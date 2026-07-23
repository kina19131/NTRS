"""
make_h1h2_pptx.py
H1 + H2 results presentation — standalone pptx in ntrs_with_lit_review style.
Run from repo root:  python3 make/make_h1h2_pptx.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette (matches ntrs_with_lit_review) ────────────────────────────────────
DARK   = RGBColor(0x1a, 0x1a, 0x2e)
BG2    = RGBColor(0x10, 0x10, 0x22)
ACCENT = RGBColor(0xe9, 0x4f, 0x37)
ACC2   = RGBColor(0x39, 0x7b, 0xbf)
LIGHT  = RGBColor(0xf5, 0xf5, 0xf5)
MID    = RGBColor(0xb0, 0xb8, 0xc8)
GREEN  = RGBColor(0x43, 0xc5, 0x9e)
YELLOW = RGBColor(0xf5, 0xc5, 0x18)
PURPLE = RGBColor(0x7c, 0x3a, 0xed)

W = Inches(13.33)
H = Inches(7.5)

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIGS = {
    "h1_density":    os.path.join(ROOT, "figures/h1/h3_density_curves_lr1e4.png"),
    "h1_bars":       os.path.join(ROOT, "figures/h1/h3_sigma_half_bars.png"),
    "h2_scatter":    os.path.join(ROOT, "figures/h2/h2_fig1_forward_scatter.png"),
    "h2_combined":   os.path.join(ROOT, "figures/h2/h2_fig2_combined_models.png"),
    "h2_traj":       os.path.join(ROOT, "figures/h2/h2_fig3_trajectories.png"),
    "h2_reversed":   os.path.join(ROOT, "figures/h2/h2_fig4_reversed.png"),
}

# ── Helpers ───────────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = DARK
    return s

def txbox(slide, text, x, y, w, h,
          size=18, bold=False, color=LIGHT,
          align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = color
    return tb

def hline(slide, x, y, w, color=ACCENT):
    ln = slide.shapes.add_connector(1, Inches(x), Inches(y), Inches(x+w), Inches(y))
    ln.line.color.rgb = color; ln.line.width = Pt(1.2)

def filled_rect(slide, x, y, w, h, fill, border=None):
    sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border
    else: sh.line.fill.background()
    return sh

def accent_bar(slide, color=ACCENT):
    filled_rect(slide, 0, 0, 0.1, 7.5, color)

def header(slide, title, subtitle=None, bar_color=ACCENT, sub_color=ACCENT):
    accent_bar(slide, bar_color)
    txbox(slide, title,    0.25, 0.12, 12.8, 0.65, size=30, bold=True)
    if subtitle:
        txbox(slide, subtitle, 0.25, 0.78, 12.8, 0.38, size=17, color=sub_color, italic=True)
    hline(slide, 0.25, 1.18, 12.8, color=bar_color)

def img(slide, key, x, y, w, h):
    path = FIGS[key]
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        sh = filled_rect(slide, x, y, w, h, BG2, MID)
        txbox(slide, f"[missing: {key}]", x+0.1, y+0.1, w-0.2, h-0.2, size=12, color=MID)

def sidebar(slide, label, body, x, y, w=3.8,
            label_color=ACC2, label_size=15, body_size=12):
    txbox(slide, label, x, y, w, 0.42, size=label_size, bold=True, color=label_color)
    filled_rect(slide, x, y+0.44, w, 0.04, BG2)  # thin divider
    txbox(slide, body,  x+0.1, y+0.52, w-0.2,
          max(0.5, len(body)*0.018), size=body_size, color=LIGHT)

def finding_box(slide, tag, title, body, x, y, w=6.3, h=1.15, tag_color=ACC2):
    sh = filled_rect(slide, x, y, w, h, BG2, tag_color)
    txbox(slide, f"{tag}  {title}", x+0.15, y+0.07, w-0.3, 0.38,
          size=14, bold=True, color=tag_color)
    txbox(slide, body,              x+0.15, y+0.48, w-0.3, h-0.55,
          size=11, color=MID)

def footer(slide, text):
    txbox(slide, text, 0.25, 6.95, 12.8, 0.45,
          size=13, color=MID, italic=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# Build slides
# ══════════════════════════════════════════════════════════════════════════════
prs = new_prs()

# ── Slide 1 — Title ──────────────────────────────────────────────────────────
s = blank(prs)
accent_bar(s, ACCENT)
txbox(s, "NTRS", 0.38, 0.45, 12.5, 1.1, size=72, bold=True, color=LIGHT)
txbox(s, "Neural Thickets  ×  Randomized Smoothing",
      0.38, 1.6, 12.0, 0.55, size=22, color=MID)
txbox(s, "H1 + H2 Experimental Results",
      0.38, 2.22, 12.0, 0.5, size=20, color=ACCENT)
hline(s, 0.38, 2.9, 5.5, color=ACC2)
txbox(s, "H1   LoRA widens the certified weight-space basin",
      0.38, 3.12, 9.5, 0.52, size=18, color=ACC2)
txbox(s, "H2   R_A = 1 predicts catastrophic sequential forgetting",
      0.38, 3.68, 9.5, 0.52, size=18, color=GREEN)
txbox(s, "Kina Kim  |  SFU  |  June 2026",
      0.38, 6.8, 9.0, 0.45, size=14, color=MID)


# ── Slide 2 — Overview table ─────────────────────────────────────────────────
s = blank(prs)
header(s, "Results at a Glance",
       "Two hypotheses tested  |  GPT-2 (124M) and GPT-2-medium (354M)")

# Column headers
for x, label, col in [(0.25,"Hypothesis",LIGHT),(4.6,"Status",LIGHT),(6.95,"Key Finding",LIGHT)]:
    txbox(s, label, x, 1.28, 2.8, 0.35, size=13, bold=True, color=col)
hline(s, 0.25, 1.65, 12.8, color=MID)

rows = [
    ("H1  LoRA widens the loss basin",
     "✅  Confirmed",
     "Pretrained σ½=0.000483  →  LoRA lr=1e-4 σ½=0.000756  (+57%)\n"
     "Higher LR (5e-4) collapses basin — training regime controls the tradeoff.\n"
     "LoRA update directions are 12% sharper than isotropic (higher curvature).",
     ACC2, GREEN),
    ("H2  R_A = 1 predicts catastrophic forgetting\n"
     "     Forward: AGNews → SST-2",
     "✅  Confirmed\n(GPT-2 + medium)",
     "R_A < 1  →  ≤ 5.6% forgetting  (safe)\n"
     "R_A > 1  →  ≥ 18% forgetting  (significant/catastrophic)\n"
     "Threshold at R_A=1 holds identically across both model scales.",
     YELLOW, GREEN),
    ("H2  Reversed: SST-2 → AGNews",
     "⚠️  Informative\nnull result",
     "Cannot reach R_A=1 with standard LRs — SST-2 σ½_A is 49× wider than AGNews.\n"
     "Finding: σ½_A encodes task specialization depth. AGNews creates a narrow basin\n"
     "(pretraining bias: WebText is news-heavy); SST-2 creates a wide, diffuse basin.",
     PURPLE, MID),
]

y = 1.75
for hyp, status, finding, hcol, scol in rows:
    h_row = 1.35
    filled_rect(s, 0.25, y, 12.85, h_row, BG2)
    txbox(s, hyp,     0.35, y+0.08, 4.15, h_row-0.2, size=11, bold=True,  color=hcol)
    txbox(s, status,  4.6,  y+0.08, 2.25, h_row-0.2, size=11, bold=True,  color=scol)
    txbox(s, finding, 6.95, y+0.05, 6.1,  h_row-0.15, size=10, bold=False, color=LIGHT)
    y += h_row + 0.08

footer(s, "σ½_A = certified radius of Task A basin  |  R_A = per-param displacement / σ½_A  "
          "|  NLL criterion for AGNews (4-class), acc-based (δ=5pp) for SST-2 (binary)")


# ── Slide 3 — H1 density curves ──────────────────────────────────────────────
s = blank(prs)
header(s, "H1 — LoRA Widens the Certified Weight-Space Basin",
       "GPT-2, WikiText-2 NLL scoring  |  lr=1e-4, rank=8  |  3-way noise direction comparison",
       bar_color=ACC2, sub_color=ACCENT)

img(s, "h1_density", 0.25, 1.28, 9.0, 5.38)

sidebar(s, "Key numbers",
        "Pretrained          σ½ = 0.000483\n"
        "LoRA isotropic  σ½ = 0.000756\n"
        "                          → +57% wider\n\n"
        "Null-space          σ½ = 0.000750\n"
        "LoRA subspace  σ½ = 0.000667\n"
        "                          → 12% narrower",
        9.45, 1.28, 3.75, label_color=GREEN)

sidebar(s, "Subspace result",
        "The LoRA update directions are\n"
        "the sharpest (highest curvature)\n"
        "— perturbations there hurt most.\n\n"
        "Null ≈ isotropic because rank=8\n"
        "is 0.00001% of d=85M params —\n"
        "geometric consequence, not a\n"
        "robustness finding.",
        9.45, 3.55, 3.75, label_color=ACC2, body_size=11)

footer(s, "Density = P[NLL(θ+ε) ≤ NLL(θ)+1e-4]  with ε~N(0,σ²I)  |  N=200 perturbations per σ  "
          "|  σ½ interpolated at density=0.5")


# ── Slide 4 — H1 bars (wider view) ───────────────────────────────────────────
s = blank(prs)
header(s, "H1 — Basin Width vs Training Regime",
       "LR=1e-4 widens; LR=5e-4 collapses — strong regime dependence",
       bar_color=ACC2, sub_color=ACCENT)

img(s, "h1_bars", 0.25, 1.28, 9.5, 5.38)

sidebar(s, "High-LR collapse",
        "At lr=5e-4: density < 0.5\n"
        "even at σ=0.0001 (smallest\n"
        "test point).\n\n"
        "Model moved so far from\n"
        "pretrained weights that even\n"
        "infinitesimal noise destroys\n"
        "performance — basin is\n"
        "smaller than measurement\n"
        "resolution.",
        9.95, 1.28, 3.2, label_color=ACCENT, body_size=11)

sidebar(s, "Implication for H2",
        "σ½_A (Task A certified radius)\n"
        "depends on which LR was used\n"
        "for Phase 1 fine-tuning.\n"
        "Conservative lr=1e-4 gives\n"
        "the most reliable σ½_A.",
        9.95, 4.5, 3.2, label_color=YELLOW, body_size=11)

footer(s, "Left bar = σ½ (×10⁻⁴)  |  Right panel: lr=5e-4 density curves for all 3 directions — none reach 0.5 threshold")


# ── Slide 5 — H2 forward scatter, side by side ───────────────────────────────
s = blank(prs)
header(s, "H2 — R_A = 1 Threshold  (Forward: AGNews → SST-2)",
       "R_A < 1  →  safe (≤5.6% forgetting)     R_A > 1  →  ≥18% forgetting",
       bar_color=GREEN, sub_color=ACCENT)

img(s, "h2_scatter", 0.25, 1.28, 9.85, 5.38)

sidebar(s, "GPT-2  (124M)\nσ½_A = 6.36×10⁻⁴",
        "R_A=0.967 → −5.6%  safe\nR_A=1.516 → −18.1%  forget\nR_A=3.017 → −31.1%",
        10.35, 1.28, 2.8, label_color=ACC2, body_size=11)

sidebar(s, "GPT-2-medium  (354M)\nσ½_A = 1.90×10⁻³",
        "R_A=0.767 → −1.1%  safe\nR_A=1.759 → −18.2%  forget\nR_A=2.923 → −50.5%\nR_A=4.889 → −49.8%",
        10.35, 3.3, 2.8, label_color=ACCENT, body_size=11)

sidebar(s, "Same threshold\nacross scales",
        "Both jump from <6%\nto >18% at R_A=1.\nσ½_A normalizes the\ncross-model comparison.",
        10.35, 5.35, 2.8, label_color=GREEN, body_size=11)

footer(s, "8 medium conditions (0620, shared Phase 1) + 6 GPT-2 conditions (0619)  |  "
          "◆ = gap-fill lr=3e-4,4e-4  |  R_A = ‖Δθ‖_RMS / σ½_A")


# ── Slide 6 — H2 combined both models ────────────────────────────────────────
s = blank(prs)
header(s, "H2 — Scale-Invariant Forgetting Budget",
       "Both models plotted on same R_A axis — threshold aligns at R_A=1 regardless of model size",
       bar_color=GREEN, sub_color=ACCENT)

img(s, "h2_combined", 0.35, 1.28, 8.7, 5.38)

txbox(s, "Why R_A works cross-model", 9.3, 1.28, 3.85, 0.42,
      size=15, bold=True, color=GREEN)
filled_rect(s, 9.3, 1.72, 3.85, 2.6, BG2)
txbox(s, "Raw displacement ‖Δθ‖ is not comparable\nacross models — medium moves 47% further\n"
         "than GPT-2 at the same LR, but forgets\nless because its basin is wider.\n\n"
         "R_A = ‖Δθ‖ / σ½_A normalizes by basin\nwidth, making both models comparable on\na single axis with the same threshold.",
      9.45, 1.78, 3.6, 2.4, size=12, color=LIGHT)

txbox(s, "Result", 9.3, 4.5, 3.85, 0.42, size=15, bold=True, color=YELLOW)
filled_rect(s, 9.3, 4.94, 3.85, 1.7, BG2)
txbox(s, "14 conditions, 2 model sizes, zero\nexceptions to R_A=1 threshold.\n\n"
         "Blue border = GPT-2  |  Red = medium\nFilled = R_A<1 (safe)  |  Warm = R_A>1",
      9.45, 5.0, 3.6, 1.55, size=12, color=LIGHT)

footer(s, "★ = medium aggressive lr=5e-4 r=32  (R_A=9.6, ΔA=−65%, from separate Phase 1 run)  "
          "|  All other points share same Phase 1 within each model")


# ── Slide 7 — H2 trajectories ────────────────────────────────────────────────
s = blank(prs)
header(s, "H2 — Forgetting Accelerates Precisely at R_A = 1",
       "Task-A accuracy during Task-B training  |  R_A grows left-to-right as Phase 2 trains",
       bar_color=GREEN, sub_color=ACCENT)

img(s, "h2_traj", 0.25, 1.28, 12.83, 5.38)

footer(s, "Left: GPT-2 — lr=1e-4 r=32 crosses R_A=1 mid-training and drops from 87% to 69%  "
          "|  Right: medium — lr=2e-4 r=32 crosses R_A=1 and drops to 72%; aggressive drops to 25%")


# ── Slide 8 — H2 reversed ────────────────────────────────────────────────────
s = blank(prs)
header(s, "H2 — Reversed (SST-2 → AGNews): A Basin Asymmetry Finding",
       "Cannot reach R_A=1 under standard conditions — SST-2 basin is 49× wider than AGNews",
       bar_color=PURPLE, sub_color=ACCENT)

img(s, "h2_reversed", 0.25, 1.28, 9.35, 5.38)

sidebar(s, "Why R_A stays << 1",
        "SST-2 acc-based σ½_A = 0.0313\n"
        "AGNews NLL-based σ½_A = 0.000636\n"
        "Ratio: 49×\n\n"
        "To reach R_A=1 in reversed\n"
        "direction: need norm > 0.031\n"
        "Current max norm = 0.00259\n"
        "→ need ~12× more displacement\n"
        "(lr ≈ 2e-3, unstable territory)",
        9.6, 1.28, 3.6, label_color=PURPLE, body_size=11)

sidebar(s, "What this means",
        "SST-2 sentiment is diffuse and\n"
        "deeply embedded in GPT-2\n"
        "(WebText pretraining encodes\n"
        "sentiment broadly).\n\n"
        "AGNews topic classification\n"
        "specializes a narrow subregion\n"
        "→ easy to overwrite.\n\n"
        "σ½_A measures task\nspecialization depth.",
        9.6, 4.0, 3.6, label_color=ACC2, body_size=11)

footer(s, "6 GPT-2 conditions, SST-2→AGNews  |  acc-based σ½_A (δ=5pp) for SST-2 Task A  "
          "|  Cannot run NLL-based here — NLL overestimates binary task basin")


# ── Slide 9 — Next Steps: H3 Rank Sweep ──────────────────────────────────────
s = blank(prs)
header(s, "Next Steps — H3: Proper Subspace Testing via Rank Sweep",
       "Current rank=8 is 0.00001% of d=85M — null-space ≈ isotropic is a geometric artifact, not a finding",
       bar_color=PURPLE, sub_color=ACCENT)

# Left: the problem
txbox(s, "The problem with rank=8", 0.25, 1.28, 6.1, 0.42,
      size=15, bold=True, color=ACCENT)
filled_rect(s, 0.25, 1.72, 6.1, 1.55, BG2)
txbox(s, "rank / d_block = 8 / 85M ≈ 0.00001%\n"
         "Any isotropic Gaussian lives almost entirely in the null-space by geometry.\n"
         "→ null ≈ isotropic is trivially guaranteed, not an empirical finding.\n"
         "→ subspace (12% sharper) is real but effect is small at this rank.",
      0.4, 1.78, 5.8, 1.42, size=11, color=LIGHT)

# Table: rank sweep plan
txbox(s, "Proposed rank sweep  (same lr=1e-4, same σ grid)", 0.25, 3.38, 8.0, 0.42,
      size=14, bold=True, color=PURPLE)

headers = [("Rank", 0.25, 0.9), ("rank/d (%)", 1.2, 1.2), ("Coverage", 2.45, 1.0),
           ("Expected signal", 3.5, 4.75)]
for label, x, w in headers:
    txbox(s, label, x, 3.82, w, 0.32, size=12, bold=True, color=MID)
hline(s, 0.25, 4.16, 8.0, color=MID)

rows_r = [
    ("8",   "0.00001%", "baseline (current)",    "null ≈ iso (trivial)"),
    ("32",  "0.00004%", "still negligible",       "null ≈ iso"),
    ("128", "0.00015%", "~0.1% coverage",         "first divergence?"),
    ("512", "0.0006%",  "still low, but larger",  "subspace vs iso gap grows"),
    ("768", "0.0009%",  "max LoRA rank for GPT-2","largest testable effect"),
]
y_r = 4.2
for rank, cov, desc, signal in rows_r:
    col = YELLOW if rank in ("512","768") else LIGHT
    for val, x, w in [(rank,0.25,0.9),(cov,1.2,1.2),(desc,2.45,1.0),(signal,3.5,4.75)]:
        txbox(s, val, x, y_r, w, 0.32, size=11, color=col)
    y_r += 0.33

# Right: hypothesis
txbox(s, "Hypothesis", 8.6, 1.28, 4.55, 0.42, size=15, bold=True, color=PURPLE)
filled_rect(s, 8.6, 1.72, 4.55, 2.2, BG2)
txbox(s, "At some rank coverage threshold,\n"
         "subspace σ½ / isotropic σ½ diverges\n"
         "significantly. Below that threshold,\n"
         "the subspace is too small to matter.\n\n"
         "Plot: subspace σ½ / iso σ½  vs  rank\n"
         "Look for the inflection point.",
      8.75, 1.78, 4.25, 2.05, size=12, color=LIGHT)

txbox(s, "Run commands", 8.6, 4.05, 4.55, 0.42, size=14, bold=True, color=GREEN)
filled_rect(s, 8.6, 4.49, 4.55, 2.2, BG2)
txbox(s, "for RANK in 8 32 128 512 768; do\n"
         "  python3 lora_density_experiment.py \\\n"
         "    --model gpt2 --lr 1e-4 \\\n"
         "    --lora_rank $RANK \\\n"
         "    --subspace --null_space \\\n"
         "    --output_dir outputs_h3_rank_sweep\n"
         "done",
      8.72, 4.55, 4.3, 2.05, size=10, color=GREEN)

footer(s, "One command per rank — lora_density_experiment.py takes single --lora_rank  |  "
          "Each run: ~20–40 min on A40  |  All runs share same σ grid for direct comparison")


# ── Slide 10 — Next Steps: Better Task Pairs for H2 ──────────────────────────
s = blank(prs)
header(s, "Next Steps — H2: Better Task Pairs for Stronger Claim",
       "AGNews/SST-2 has two weaknesses — pretraining bias asymmetry + 49× basin width mismatch",
       bar_color=GREEN, sub_color=ACCENT)

# Problems with current pair
txbox(s, "Problems with AGNews → SST-2", 0.25, 1.28, 6.1, 0.42,
      size=14, bold=True, color=ACCENT)
filled_rect(s, 0.25, 1.72, 6.1, 1.35, BG2)
txbox(s, "①  Pretraining bias: GPT-2 trained on WebText (news-heavy) → AGNews is 'easy to specialize.'\n"
         "②  Basin asymmetry: SST-2 σ½_A is 49× wider than AGNews — reversed direction can't reach R_A=1.\n"
         "③  Weak reversed experiment limits H2 to one direction only.",
      0.4, 1.78, 5.8, 1.2, size=11, color=LIGHT)

# Recommended pairs table
txbox(s, "Recommended pairs  (all available in basin_widening_experiment.py)", 0.25, 3.15, 9.5, 0.42,
      size=14, bold=True, color=GREEN)

col_hdrs = [("Task pair", 0.25, 3.5), ("Type", 3.8, 2.0),
            ("Why better", 5.85, 4.0), ("Priority", 9.9, 2.9)]
for label, x, w in col_hdrs:
    txbox(s, label, x, 3.6, w, 0.3, size=12, bold=True, color=MID)
hline(s, 0.25, 3.92, 12.8, color=MID)

pairs = [
    ("AGNews → DBpedia",
     "Topic → Topic\n(4-class → 14-class)",
     "Both topic classification, no NLI/sentiment bias.\nMore symmetric basins. Pretraining neutral.",
     "★ Try first\n(easiest)"),
    ("DBpedia → AGNews",
     "Reversed of above",
     "Tests whether threshold holds in both directions.\nDBpedia is fine-grained, AGNews is coarse — directional.",
     "★ Run together"),
    ("MNLI → AGNews",
     "NLI → Topic\n(3-class → 4-class)",
     "Structurally different tasks, no pretraining bias for either.\n"
     "MNLI fine-tuning doesn't specialize a narrow news-basin.",
     "Good backup"),
]
y_p = 3.98
for pair, typ, why, pri in pairs:
    h_p = 0.82
    filled_rect(s, 0.25, y_p, 12.8, h_p, BG2)
    txbox(s, pair, 0.35, y_p+0.05, 3.4, h_p-0.1, size=11, bold=True, color=YELLOW)
    txbox(s, typ,  3.8,  y_p+0.05, 1.95, h_p-0.1, size=10, color=MID)
    txbox(s, why,  5.85, y_p+0.05, 3.95, h_p-0.1, size=10, color=LIGHT)
    txbox(s, pri,  9.9,  y_p+0.05, 2.8, h_p-0.1, size=11, bold=True, color=GREEN)
    y_p += h_p + 0.05

txbox(s, "Key criterion: both tasks >80% accuracy on GPT-2  |  σ½_A within ~5× in both directions  "
         "|  Neither task has large pretraining advantage",
      0.25, 6.6, 12.8, 0.38, size=12, color=MID, italic=True)

footer(s, "RTE/QQP/MRPC not in current script — would need to add loaders  |  "
          "DBpedia: 14-class ontology  |  MNLI: 3-class NLI (entail/neutral/contradict)")


# ── save ─────────────────────────────────────────────────────────────────────
out = os.path.join(ROOT, "presentations", "ntrs_h1h2_results.pptx")
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
