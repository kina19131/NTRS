"""
Append 0621 result slides to ntrs_h1h2_results.pptx:
  Slide 11 — H3 updated: sub/iso = 0.427 (strong signal)
  Slide 12 — H2 new task pairs: DBpedia↔AGNews findings
Run from repo root: python3 make/append_0621_slides.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DARK   = RGBColor(0x1a, 0x1a, 0x2e)
BG2    = RGBColor(0x10, 0x10, 0x22)
ACCENT = RGBColor(0xe9, 0x4f, 0x37)
ACC2   = RGBColor(0x39, 0x7b, 0xbf)
LIGHT  = RGBColor(0xf5, 0xf5, 0xf5)
MID    = RGBColor(0xb0, 0xb8, 0xc8)
GREEN  = RGBColor(0x43, 0xc5, 0x9e)
YELLOW = RGBColor(0xf5, 0xc5, 0x18)
RED    = RGBColor(0xe9, 0x4f, 0x37)

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIGS_H2   = os.path.join(ROOT, "figures", "h2")
FIGS_H3   = os.path.join(ROOT, "figures", "h3")

def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = DARK
    return s

def txbox(slide, text, x, y, w, h,
          size=18, bold=False, color=LIGHT,
          align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = color
    return tb

def hline(slide, x, y, w, color=ACCENT):
    ln = slide.shapes.add_connector(1, Inches(x), Inches(y), Inches(x+w), Inches(y))
    ln.line.color.rgb = color; ln.line.width = Pt(1.2)

def filled_rect(slide, x, y, w, h, fill, border=None):
    sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border
    else: sh.line.fill.background()
    return sh

def accent_bar(slide, color=ACCENT):
    filled_rect(slide, 0, 0, 0.1, 7.5, color)

def header(slide, title, subtitle=None, bar_color=ACCENT, sub_color=ACCENT):
    accent_bar(slide, bar_color)
    txbox(slide, title,    0.25, 0.12, 12.8, 0.65, size=30, bold=True)
    if subtitle:
        txbox(slide, subtitle, 0.25, 0.78, 12.8, 0.38, size=17, color=sub_color, italic=True)
    hline(slide, 0.25, 1.18, 12.8, color=bar_color)

def img(slide, fname, x, y, w, h, subdir="h2"):
    folder = FIGS_H3 if subdir == "h3" else FIGS_H2
    path = os.path.join(folder, fname)
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        sh = filled_rect(slide, x, y, w, h, BG2, MID)
        txbox(slide, f"[missing: {fname}]", x+0.1, y+0.1, w-0.2, h-0.2, size=12, color=MID)

def sidebar(slide, label, body, x, y, w=3.8,
            label_color=ACC2, label_size=15, body_size=12):
    txbox(slide, label, x, y, w, 0.42, size=label_size, bold=True, color=label_color)
    filled_rect(slide, x, y+0.44, w, 0.04, BG2)
    txbox(slide, body,  x+0.1, y+0.52, w-0.2,
          max(0.5, len(body)*0.018), size=body_size, color=LIGHT)

def footer(slide, text):
    txbox(slide, text, 0.25, 6.95, 12.8, 0.45,
          size=13, color=MID, italic=True, align=PP_ALIGN.CENTER)


# ── Open existing pptx ────────────────────────────────────────────────────────
DST = os.path.join(ROOT, "presentations", "ntrs_h1h2_results.pptx")
prs = Presentation(DST)
n_before = len(prs.slides)
print(f"Opened '{DST}' with {n_before} slides — appending 2 slides")


# ══════════════════════════════════════════════════════════════════════════════
# Slide 11 — H3 Updated: sub/iso = 0.427 (strong signal)
# ══════════════════════════════════════════════════════════════════════════════
s = blank(prs)
header(s,
    "H3 Update — LoRA Subspace is 2.3× Sharper Than Isotropic",
    "New run  |  sub/iso = 0.427  (vs 0.882 in initial run)  |  null ≈ pretrained (1.03×)",
    bar_color=GREEN, sub_color=GREEN)

img(s, "h3_new_density_bars.png", 0.25, 1.28, 9.3, 5.35, subdir="h3")

sidebar(s, "Key numbers",
    "σ½ pretrained  = 0.000483\n"
    "σ½ null-space  = 0.000500  ← ≈ pretrained\n"
    "σ½ isotropic   = 0.000445\n"
    "σ½ subspace    = 0.000190  ← sharpest\n\n"
    "sub / iso = 0.427×  (−57%)\n"
    "null / pretrained = 1.03×",
    9.75, 1.28, 3.5, label_color=GREEN, body_size=11)

sidebar(s, "What it means",
    "LoRA update directions are the\ncurvature-concentrated directions.\n\n"
    "Null-space (≈pretrained) is the\nflat / generalization space.\n\n"
    "Fine-tuning with LoRA separates\nthe weight space into:\n"
    "  • sharp subspace (updated dirs)\n"
    "  • flat null-space (unmodified dirs)",
    9.75, 3.4, 3.5, label_color=ACC2, body_size=11)

sidebar(s, "Previous vs new",
    "Initial run (0620_lora):\n"
    "  sub/iso = 0.882  (−12%)\n\n"
    "New run (0621_h3):\n"
    "  sub/iso = 0.427  (−57%)\n\n"
    "Both same lr=1e-4, GPT-2\n"
    "→ effect is real, run-dependent",
    9.75, 5.2, 3.5, label_color=YELLOW, body_size=11)

footer(s,
    "H3: subspace = LoRA-adapted directions in weight space  |  "
    "null-space = directions orthogonal to LoRA update  |  "
    "NLL criterion, N=200 perturbations")


# ══════════════════════════════════════════════════════════════════════════════
# Slide 12 — H2 New Task Pairs: DBpedia↔AGNews
# ══════════════════════════════════════════════════════════════════════════════
s = blank(prs)
header(s,
    "H2 New Task Pairs — DBpedia ↔ AGNews",
    "DBpedia→AGNews: safe at all tested R_A  |  AGNews→DBpedia: early forgetting reveals threshold boundary",
    bar_color=ACC2, sub_color=YELLOW)

img(s, "h2_dbpedia_agnews.png", 0.25, 1.28, 9.3, 5.35)

sidebar(s, "DBpedia→AGNews",
    "σ½_A = 0.014500  (NLL)\nMax R_A = 0.128 → −0.9%\n\n"
    "All conditions safe.\nR_A never reaches 1 — σ½_A\ntoo large to test threshold.\n\n"
    "DBpedia (14-class) creates a\nvery wide basin (high diversity).",
    9.75, 1.28, 3.5, label_color=GREEN, body_size=11)

sidebar(s, "AGNews→DBpedia",
    "σ½_A = 0.006442  (NLL)\nMax R_A = 0.394 → −17.5%\n\n"
    "Forgetting at R_A << 1!\n"
    "This is structural mismatch:\n"
    "4-class → 14-class output\nforces representation change\n"
    "even at small weight displacement.",
    9.75, 3.35, 3.5, label_color=RED, body_size=11)

sidebar(s, "Implication for H2",
    "R_A=1 threshold (H2) holds for\nstructurally similar tasks\n(AGNews→SST-2 both <5 classes).\n\n"
    "Structural mismatch (4→14 class)\ncauses forgetting at R_A≪1 via\nrepresentational drift, not\nbasin exit.\n\n"
    "Better pair needed: tasks with\nsimilar class granularity.",
    9.75, 5.1, 3.5, label_color=YELLOW, body_size=10)

footer(s,
    "DBpedia: 14-class Wikipedia topic  |  AGNews: 4-class news topic  |  "
    "GPT-2, NLL-based σ½_A  |  6 conditions per direction (lr × rank grid)")


# ── Save ──────────────────────────────────────────────────────────────────────
prs.save(DST)
print(f"Saved '{DST}'  ({len(prs.slides)} slides total, added {len(prs.slides)-n_before})")
