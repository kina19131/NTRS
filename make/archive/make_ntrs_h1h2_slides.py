"""
Append H1 + H2 results slides to ntrs_presentation3.pptx.
Matches the existing dark-theme slide format exactly.
"""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import shutil, os

# ── colour constants (from existing slides) ───────────────────────────────────
C_WHITE   = RGBColor(0xF5, 0xF5, 0xF5)
C_RED     = RGBColor(0xE9, 0x4F, 0x37)
C_BLUE    = RGBColor(0x39, 0x7B, 0xBF)
C_TEAL    = RGBColor(0x43, 0xC5, 0x9E)
C_YELLOW  = RGBColor(0xF5, 0xC5, 0x18)
C_GRAY    = RGBColor(0xB0, 0xB8, 0xC8)
C_DARK    = RGBColor(0x1A, 0x1A, 0x2E)
C_BAR     = RGBColor(0x39, 0x7B, 0xBF)

IN = Inches

def rgb_fill(shape, color: RGBColor):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, text, left, top, width, height,
                font_size=12, bold=False, color=C_WHITE,
                align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(IN(left), IN(top), IN(width), IN(height))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def add_separator(slide, left, top, width):
    cx = IN(width); cy = Emu(0)
    connector = slide.shapes.add_connector(
        1, IN(left), IN(top), IN(left) + cx, IN(top))
    connector.line.color.rgb = RGBColor(0x44, 0x44, 0x55)
    connector.line.width = Pt(0.75)

def add_image(slide, path, left, top, width, height):
    return slide.shapes.add_picture(path, IN(left), IN(top), IN(width), IN(height))

def add_filled_box(slide, left, top, width, height, fill_color, alpha=None):
    box = slide.shapes.add_shape(1, IN(left), IN(top), IN(width), IN(height))
    rgb_fill(box, fill_color)
    box.line.fill.background()
    return box

def add_slide(prs):
    """Add blank slide with dark background matching the theme."""
    blank_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank_layout)
    # Dark background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x0F, 0x0F, 0x1A)
    return slide

def std_header(slide, title, subtitle, bar_color=C_BAR):
    """Standard slide header: thin left bar + title + subtitle + separator."""
    # Left accent bar
    bar = add_filled_box(slide, 0, 0, 0.1, 7.5, bar_color)
    # Title
    add_textbox(slide, title, 0.25, 0.12, 12.8, 0.62,
                font_size=32, bold=True, color=C_WHITE)
    # Subtitle
    add_textbox(slide, subtitle, 0.25, 0.75, 12.8, 0.38,
                font_size=18, bold=False, color=C_RED)
    # Separator
    add_separator(slide, 0.25, 1.18, 12.8)

def sidebar_block(slide, label, body, x, y, w,
                  label_color=C_TEAL, label_size=17, body_size=13):
    add_textbox(slide, label, x, y, w, 0.4,
                font_size=label_size, bold=True, color=label_color)
    box = add_filled_box(slide, x, y+0.44, w, 0.04+len(body)*0.017,
                         RGBColor(0x1E, 0x1E, 0x30))
    body_h = max(0.5, len(body) * 0.185 + 0.2)
    add_textbox(slide, body, x+0.15, y+0.5, w-0.3, body_h,
                font_size=body_size, bold=False, color=C_WHITE)

def std_footer(slide, text):
    add_textbox(slide, text, 0.25, 6.95, 12.8, 0.45,
                font_size=13, bold=False, color=C_GRAY)


# ── open existing pptx ────────────────────────────────────────────────────────
src = "ntrs_presentation3.pptx"
dst = "ntrs_presentation3.pptx"
prs = Presentation(src)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE A — Overview table: H1 + H2 summary
# ════════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs)
std_header(slide,
    "H1 + H2 Results Summary",
    "Sequential Forgetting Budget + LoRA Basin Geometry  |  GPT-2 & GPT-2-medium")

# Column headers
for x, label, col in [(0.25, "Hypothesis", C_WHITE),
                       (4.5,  "Status",     C_WHITE),
                       (6.8,  "Key Finding", C_WHITE)]:
    add_textbox(slide, label, x, 1.28, 3.8, 0.38,
                font_size=14, bold=True, color=col)
add_separator(slide, 0.25, 1.68, 12.8)

rows = [
    ("H1  LoRA widens the loss basin",
     "✅  Confirmed",
     "σ½: 0.000483 (pretrained) → 0.000756 (LoRA lr=1e-4)  |  +57%\n"
     "Higher LR collapses basin — training regime determines effect",
     C_TEAL),
    ("H2  R_A = 1 predicts catastrophic forgetting",
     "✅  Confirmed\n(GPT-2 + medium)",
     "R_A < 1 → ≤5.6% forgetting  |  R_A > 1 → ≥18% forgetting\n"
     "Threshold holds across model scale (124M and 354M params)\n"
     "Trajectory shows forgetting accelerates exactly at R_A = 1",
     C_YELLOW),
]

y = 1.78
for hyp, status, finding, col in rows:
    h = 1.4
    add_filled_box(slide, 0.25, y, 13.05, h, RGBColor(0x1A, 0x1A, 0x2E))
    add_textbox(slide, hyp,     0.35, y+0.08, 4.0, h-0.2,
                font_size=13, bold=True, color=col)
    add_textbox(slide, status,  4.5,  y+0.08, 2.2, h-0.2,
                font_size=13, bold=True, color=C_TEAL)
    add_textbox(slide, finding, 6.85, y+0.05, 6.3, h-0.1,
                font_size=11, bold=False, color=C_WHITE)
    y += h + 0.1

# Note on H3
add_textbox(slide,
    "H3  LoRA subspace certification  —  Partial / Revised: "
    "Subspace σ½ (0.000667) is 12% narrower than isotropic (0.000756). "
    "Update subspace is the sharpest direction; null-space ≈ isotropic "
    "is a geometric consequence of rank (8) << d (85M), not an independent finding.",
    0.25, 5.2, 12.8, 1.1,
    font_size=12, bold=False, color=C_GRAY)

std_footer(slide,
    "H2: AGNews→SST-2 forward  |  σ½_A measured with NLL (GPT-2 σ½=6.36e-4, Medium σ½=1.90e-3)  "
    "|  R_A = per-param-norm / σ½_A")


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE B — H1: LoRA Basin Widening
# ════════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs)
std_header(slide,
    "H1 — LoRA Widens the Loss Basin",
    "Isotropic noise density  |  GPT-2, WikiText-2 NLL  |  lr=1e-4, rank=8")

add_image(slide, "outputs_0620_lora/figures/h3_density_curves_lr1e4.png",
          0.25, 1.28, 8.9, 5.35)

sidebar_block(slide, "Key result",
    "Pretrained   σ½ = 0.000483\nLoRA lr=1e-4  σ½ = 0.000756\n                    (+57%  wider)",
    9.35, 1.28, 3.75, label_color=C_TEAL)

sidebar_block(slide, "What this means",
    "After LoRA fine-tuning, random\n"
    "weight perturbations up to 57%\n"
    "larger are tolerated without loss\n"
    "degradation.\n\n"
    "Higher LR (5e-4) collapses the\n"
    "basin — all σ start below 0.5.\n"
    "Training regime controls the tradeoff.",
    9.35, 3.25, 3.75, label_color=C_BLUE, body_size=12)

sidebar_block(slide, "H3 note",
    "Subspace σ½ = 0.000667 (<iso)\n"
    "→ LoRA update dirs are sharpest\n"
    "Null ≈ iso: geometric artifact\n"
    "(rank 8 << d=85M)",
    9.35, 5.45, 3.75, label_color=C_GRAY, body_size=11)

std_footer(slide,
    "Density = fraction of N=200 perturbations where NLL(θ+ε) ≤ NLL(θ) + 1e-4  |  "
    "σ½ = largest σ with density ≥ 0.5 (interpolated)")


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE C — H2 Forward Scatter: GPT-2 + Medium
# ════════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs)
std_header(slide,
    "H2 — R_A = 1 Threshold  (Forward: AGNews → SST-2)",
    "R_A < 1 → safe (≤5.6% forgetting)   |   R_A > 1 → catastrophic forgetting (≥18%)")

add_image(slide, "outputs_0620_h2/figures/h2_fig1_forward_scatter.png",
          0.25, 1.28, 9.7, 5.35)

sidebar_block(slide, "GPT-2 (124M)",
    "σ½_A = 0.000636  (NLL)\nMax safe R_A:  0.967 → −5.6%\nFirst forget:  1.516 → −18.1%\nAt R_A=3.0:        → −31.1%",
    10.15, 1.28, 3.05, label_color=C_BLUE, body_size=11)

sidebar_block(slide, "GPT-2-medium (354M)",
    "σ½_A = 0.001900  (NLL)\nMax safe R_A:  0.767 → −1.1%\nFirst forget:  1.759 → −18.2%\nAt R_A=4.9:        → −50%",
    10.15, 3.2, 3.05, label_color=C_RED, body_size=11)

sidebar_block(slide, "Scale result",
    "Both models show same\n"
    "threshold at R_A = 1.\n"
    "R_A normalizes cross-\n"
    "model comparisons.",
    10.15, 5.15, 3.05, label_color=C_TEAL, body_size=11)

std_footer(slide,
    "8 medium conditions (0620, same Phase 1) + 6 GPT-2 conditions (0619)  |  "
    "◆ = gap-fill conditions lr=3e-4, 4e-4  |  R_A = ‖θ_B−θ_A‖_per-param / σ½_A")


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE D — H2 Combined + Trajectories
# ════════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs)
std_header(slide,
    "H2 — R_A = 1 Generalizes Across Scale + Trajectory Evidence",
    "Both models align at R_A=1  |  Forgetting accelerates exactly when R_A crosses 1 during training")

# Two figures side by side
add_image(slide, "outputs_0620_h2/figures/h2_fig2_combined_models.png",
          0.25, 1.28, 6.5, 5.35)
add_image(slide, "outputs_0620_h2/figures/h2_fig3_trajectories.png",
          6.85, 1.28, 6.35, 5.35)

std_footer(slide,
    "Left: all 14 conditions, both models, one plot — blue border=GPT-2, red border=medium  |  "
    "Right: Task-A accuracy during Task-B training — R_A grows left to right as Phase 2 trains")


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE E — H2 Reversed + Interpretation
# ════════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs)
std_header(slide,
    "H2 — Reversed (SST-2 → AGNews): Basin Asymmetry Finding",
    "Cannot reach R_A=1 with standard LRs — SST-2 basin 49× wider than AGNews")

add_image(slide, "outputs_0620_h2/figures/h2_fig4_reversed.png",
          0.25, 1.28, 9.2, 5.35)

sidebar_block(slide, "What happened",
    "SST-2 σ½_A = 0.031250\n(acc-based, δ=5pp)\nMax norm at lr=2e-4:\n  = 0.00259\n→ R_A_max = 0.083\n\nCannot reach R_A=1\nwith normal fine-tuning.",
    9.6, 1.28, 3.7, label_color=C_RED, body_size=11)

sidebar_block(slide, "What it means",
    "SST-2's basin is 49× wider\nthan AGNews's. GPT-2's\nWebText pretraining is\nnews-heavy → AGNews\nspecializes sharply (narrow\nbasin), SST-2 sentiment is\ndiffuse (wide basin).\n\nσ½_A predicts WHICH task\nsequences are at risk —\nnot just when forgetting\noccurs.",
    9.6, 3.7, 3.7, label_color=C_BLUE, body_size=11)

std_footer(slide,
    "Reversed: SST-2 as Task A (acc-based σ½), AGNews as Task B  |  "
    "6 conditions, GPT-2  |  All R_A < 0.09 — far from threshold")

# ── save ─────────────────────────────────────────────────────────────────────
prs.save(dst)
print(f"Saved: {dst}  ({len(prs.slides)} slides total)")
