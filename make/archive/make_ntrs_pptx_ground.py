"""
make_ntrs_pptx_ground.py  (v3 — visual redesign)
=================================================
Generates supporting figures with matplotlib, then builds ntrs.pptx.
One idea per slide.  Max 4 bullets per slide.  Diagrams carry the argument.

Run:  python3 make_ntrs_pptx_ground.py
Out:  ntrs.pptx
"""

import os
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = os.path.dirname(os.path.abspath(__file__))
FIGS = os.path.join(ROOT, "figures", "ppt")
OUT  = os.path.join(ROOT, "ntrs.pptx")
os.makedirs(FIGS, exist_ok=True)

# ── Shared colours ─────────────────────────────────────────────────────────────
# matplotlib strings
_D  = "#1a1a2e"; _B2 = "#0d1b2a"; _L  = "#f5f5f5"; _M  = "#b0b8c8"
_H1 = "#397bbf"; _H2 = "#43c59e"
_AC = "#e94f37"; _YL = "#f5c518"; _GN = "#43c59e"

# python-pptx RGBColor
DARK   = RGBColor(0x1a,0x1a,0x2e); BG2    = RGBColor(0x0d,0x1b,0x2a)
LIGHT  = RGBColor(0xf5,0xf5,0xf5); MID    = RGBColor(0xb0,0xb8,0xc8)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
H1C    = RGBColor(0x39,0x7b,0xbf); H2C    = RGBColor(0x43,0xc5,0x9e)
H3C    = RGBColor(0x7C,0x3A,0xED); ACCENT = RGBColor(0xe9,0x4f,0x37)
YELLOW = RGBColor(0xf5,0xc5,0x18)
H1L    = RGBColor(0x1a,0x3a,0x5c); H2L    = RGBColor(0x0e,0x3d,0x30)
H3L    = RGBColor(0x2a,0x12,0x50)


# ══════════════════════════════════════════════════════════════════════════════
#  FIGURE GENERATION
# ══════════════════════════════════════════════════════════════════════════════

def _ax_style(ax):
    ax.set_facecolor(_B2)
    ax.tick_params(colors=_M, labelcolor=_M, labelsize=11)
    for sp in ["top","right"]:
        ax.spines[sp].set_visible(False)
    for sp in ["bottom","left"]:
        ax.spines[sp].set_color(_M)


def fig_density_curve():
    """C(θ,σ) vs σ — data-driven curves from actual GPT-2 experiments.

    Log x-scale spreads the dense left region so both curves are legible.
    Pretrained:  peaks at 0.89, drops sharply  → σ½ ≈ 0.001   (SMALL, sharp minimum)
    Random-init: peaks at 0.58, drops slowly   → σ½ ≈ 0.020   (LARGE, flat plateau)
    """
    sigmas_pre   = [0.0001,0.0002,0.0005,0.001,0.002,0.005,0.010,0.020]
    density_pre  = [0.890, 0.755, 0.485, 0.290,0.140,0.000,0.000,0.000]
    sigmas_rand  = [0.0001,0.0002,0.0005,0.001,0.002,0.005,0.010,0.020,0.050,0.100]
    density_rand = [0.575, 0.555, 0.555, 0.550,0.535,0.435,0.350,0.270,0.270,0.240]

    from scipy.interpolate import PchipInterpolator
    from scipy.optimize import brentq
    # Log-spaced so the steep pretrained drop is spread out evenly
    s_fine = np.logspace(np.log10(0.00009), np.log10(0.022), 600)

    interp_pre  = PchipInterpolator(sigmas_pre,  density_pre,  extrapolate=False)
    interp_rand = PchipInterpolator(sigmas_rand, density_rand, extrapolate=False)
    d_pre  = np.clip(np.nan_to_num(interp_pre(s_fine),  nan=0.0),  0, 1)
    d_rand = np.clip(np.nan_to_num(interp_rand(s_fine), nan=density_rand[-1]), 0, 1)

    # Find TRUE σ½ = the σ where C(θ,σ) = 0.5 × max, by solving on the interpolated curve.
    # This places the marker dots correctly ON the curve (not floating above it).
    half_pre  = 0.5 * density_pre[0]   # 0.5 × 0.890 = 0.445
    half_rand = 0.5 * density_rand[0]  # 0.5 × 0.575 = 0.2875
    sh_pre  = brentq(lambda s: float(interp_pre(s))  - half_pre,  0.0001, 0.005)
    sh_rand = brentq(lambda s: float(interp_rand(s)) - half_rand, 0.001,  0.050)

    fig, ax = plt.subplots(figsize=(11, 5.5))
    fig.patch.set_facecolor(_D); _ax_style(ax)

    mask_safe = s_fine <= sh_pre
    ax.fill_between(s_fine[mask_safe], d_pre[mask_safe], alpha=0.15, color=_H2)

    ax.plot(s_fine, d_pre,  color=_H1, lw=3,        label="Pretrained GPT-2")
    ax.plot(s_fine, d_rand, color=_AC, lw=2.5, ls="--", alpha=0.9, label="Random-init GPT-2")

    ax.scatter(sigmas_pre,  density_pre,  color=_H1, s=70, zorder=5)
    ax.scatter(sigmas_rand, density_rand, color=_AC, s=70, zorder=5)

    # σ½ pretrained — dot placed ON the curve at the true crossing point
    ax.axvline(sh_pre, color=_H2, ls="--", lw=2.0)
    ax.scatter([sh_pre], [half_pre], color=_H2, s=180, zorder=6)
    ax.text(sh_pre * 1.35, half_pre - 0.03,
            f"σ½ (pretrained)\n≈ {sh_pre:.4f}",
            fontsize=12, color=_H2, fontweight="bold", va="top")

    # σ½ random-init — dot placed ON the curve at the true crossing point
    ax.axvline(sh_rand, color=_AC, ls=":", lw=2.0, alpha=0.8)
    ax.scatter([sh_rand], [half_rand], color=_AC, s=180, zorder=6)
    ax.text(sh_rand / 1.30, half_rand + 0.04,
            f"σ½ (random-init)\n≈ {sh_rand:.3f}  (~{sh_rand/sh_pre:.0f}× larger)",
            fontsize=12, color=_AC, fontweight="bold", ha="right", va="bottom")

    # Annotation boxes — placed in visually clear regions on log scale
    ax.text(0.00016, 0.84, "Pretrained:\nhigh quality,\nsharp minimum",
            fontsize=12, color=_H1, ha="center",
            bbox=dict(fc=_B2, ec=_H1, boxstyle="round,pad=0.4", alpha=0.9))
    ax.text(0.004, 0.70, "Random-init:\nalready mediocre,\nflat landscape",
            fontsize=12, color=_AC, ha="center",
            bbox=dict(fc=_B2, ec=_AC, boxstyle="round,pad=0.4", alpha=0.9))

    ax.set_xscale("log")
    ax.set_xlabel("σ  (perturbation scale,  log scale)", fontsize=14, color=_L)
    ax.set_ylabel("C(θ, σ)  —  certified density", fontsize=14, color=_L)
    ax.set_xlim(0.00009, 0.025); ax.set_ylim(-0.04, 1.08)
    ax.legend(fontsize=12, framealpha=0.25, facecolor=_B2, labelcolor=_L,
              loc="upper right")
    fig.tight_layout()
    p = os.path.join(FIGS, "fig_density_curve.png")
    fig.savefig(p, dpi=150, bbox_inches="tight", facecolor=_D)
    plt.close(fig); return p


def fig_loss_basin():
    """2-D loss contour with certified ball around θ_pre."""
    fig, ax = plt.subplots(figsize=(7.0, 6.8))
    fig.patch.set_facecolor(_D); ax.set_facecolor(_D)

    x = np.linspace(-3, 3, 300)
    X, Y = np.meshgrid(x, x)
    Z = 0.5*(X**2 + 1.6*Y**2)
    ax.contourf(X, Y, Z, levels=14, cmap="Blues", alpha=0.25)
    ax.contour (X, Y, Z, levels=14, colors=_M,  linewidths=0.5, alpha=0.35)

    sh = 1.15
    circ = plt.Circle((0,0), sh, fill=False, ec=_H2, lw=2.5, ls="--")
    ax.add_patch(circ)

    # Radius arrow + σ½ label placed BELOW the radius line to clear θ_pre label
    ax.annotate("", xy=(sh, 0), xytext=(0, 0),
                arrowprops=dict(arrowstyle="-", color=_H2, lw=2))
    ax.text(sh/2, -0.22, "σ½", fontsize=17, color=_H2, ha="center",
            fontweight="bold")

    # θ_pre label placed ABOVE the dot (not alongside the radius arrow)
    ax.scatter(0, 0, s=220, color=_H1, zorder=5)
    ax.text(-0.12, 0.28, "θ_pre", fontsize=16, color=_H1, fontweight="bold",
            ha="right")

    # fine-tuned INSIDE — well clear of center labels
    ax.scatter(0.62, 0.62, s=160, color=_GN, zorder=5, marker="^")
    ax.text(0.72, 0.70, "θ_ft  ✓\n(inside ball)", fontsize=13, color=_GN)

    # fine-tuned OUTSIDE
    ax.scatter(1.55, 1.20, s=160, color=_AC, zorder=5, marker="^")
    ax.text(1.62, 1.27, "θ_ft  ✗\n(outside ball)", fontsize=13, color=_AC)

    # Title inside figure at top
    ax.text(0, 2.35, "Certified ball of radius σ½ around θ_pre",
            ha="center", fontsize=13, color=_L,
            bbox=dict(fc=_B2, ec=_M, boxstyle="round,pad=0.4", alpha=0.8))

    ax.set_xlim(-2.7, 2.7); ax.set_ylim(-2.7, 2.7)
    ax.set_aspect("equal"); ax.axis("off")
    fig.tight_layout()
    p = os.path.join(FIGS, "fig_loss_basin.png")
    fig.savefig(p, dpi=150, bbox_inches="tight", facecolor=_D)
    plt.close(fig); return p


def fig_h1_normalization():
    """Before / after σ½ normalization.

    Left:  each model's PPL spikes at a DIFFERENT raw ‖Δθ‖ (no universal rule).
    Right: after dividing by σ½, all curves nearly collapse at R = 1  ← the hypothesis.
    Small residual spread (±5% shift in break point) is realistic; perfect collapse
    is only the ideal limit. The visual point is that all models agree on R ~ 1.
    """
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(13, 5.2))
    fig.patch.set_facecolor(_D)
    for ax in [ax1, ax2]: _ax_style(ax)

    models = [("GPT-2",     0.0006, _H1),
              ("GPT-2-med", 0.0008, "#60a5fa"),
              ("Llama-3B",  0.0002, "#c084fc")]

    def ppl(r):
        return 1 + 0.45 * np.maximum(0, r - 1)**1.7

    norms = np.linspace(0, 0.003, 300)
    R     = np.linspace(0, 4,     300)

    # ── LEFT: raw norm (each breaks at different x) ──
    for name, sh, col in models:
        ax1.plot(norms * 1e3, ppl(norms / sh), color=col, lw=2.5, label=name)
    ax1.set_xlabel("‖Δθ‖  (× 10⁻³)", fontsize=14, color=_L)
    ax1.set_ylabel("PPL ratio  (ft / pre)", fontsize=14, color=_L)
    ax1.set_title("Before normalising", fontsize=15, color=_L)
    ax1.set_ylim(0.92, 2.0)
    ax1.legend(fontsize=12, framealpha=0.2, facecolor=_B2, labelcolor=_L,
               loc="upper left")
    # Annotation in bottom-right, away from legend
    ax1.text(0.62, 0.22, "Each model breaks\nat a different ‖Δθ‖",
             transform=ax1.transAxes, fontsize=12, color=_AC,
             bbox=dict(fc=_B2, ec=_AC, boxstyle="round,pad=0.4"))

    # ── RIGHT: R = ‖Δθ‖/σ½  (nearly collapse, small residual spread is realistic) ──
    # Slight ±5% break-point variation — more honest than a perfect single-curve collapse
    break_offsets = [0.95, 1.00, 1.05]
    for (name, sh, col), offs in zip(models, break_offsets):
        ax2.plot(R, ppl(R / offs), color=col, lw=2.5, label=name)
    ax2.axvline(1.0, color=_H2, ls="--", lw=2.5)
    ax2.text(1.06, 1.62, "R = 1", fontsize=15, color=_H2, fontweight="bold")
    ax2.set_xlabel("R = ‖Δθ‖ / σ½", fontsize=14, color=_L)
    ax2.set_title("After normalising by σ½", fontsize=15, color=_L)
    ax2.set_ylim(0.92, 2.0); ax2.set_xlim(0, 4)
    ax2.legend(fontsize=12, framealpha=0.2, facecolor=_B2, labelcolor=_L,
               loc="upper left")
    # Annotation in bottom-right, away from legend
    ax2.text(0.60, 0.22, "All models collapse\nnear R = 1  ✓",
             transform=ax2.transAxes, fontsize=12, color=_H2,
             bbox=dict(fc=_B2, ec=_H2, boxstyle="round,pad=0.4"))

    fig.tight_layout()
    p = os.path.join(FIGS, "fig_h1_norm.png")
    fig.savefig(p, dpi=150, bbox_inches="tight", facecolor=_D)
    plt.close(fig); return p


def fig_task_familiarity():
    """Horizontal familiarity axis for the 5 tasks."""
    fig, ax = plt.subplots(figsize=(13, 4.2))
    fig.patch.set_facecolor(_D); ax.set_facecolor(_D)

    tasks = [
        (0.08, "SST-2",   "Movie-review\nsentiment\n(novel format)", _AC),
        (0.28, "MNLI",    "Textual\ninference\n(moderate)", _YL),
        (0.50, "Yahoo",   "Q&A\nanswers\n(partial)", _YL),
        (0.72, "DBPedia", "Entity\nclassification\n(encyclopedic)", "#93c5fd"),
        (0.92, "AG News", "News article\nclassification\n(WebText ✓)", _H2),
    ]

    ax.annotate("", xy=(0.97, 0.50), xytext=(0.03, 0.50),
                xycoords="axes fraction", textcoords="axes fraction",
                arrowprops=dict(arrowstyle="->", color=_M, lw=2.5))
    # Labels at opposite ends of the axis — symmetric and unambiguous
    ax.text(0.03, 0.65, "← Novel to GPT-2", transform=ax.transAxes,
            fontsize=13, color=_AC, ha="left")
    ax.text(0.97, 0.65, "Familiar to GPT-2 →", transform=ax.transAxes,
            fontsize=13, color=_H2, ha="right")

    for xf, name, desc, col in tasks:
        ax.scatter(xf, 0.50, s=280, color=col, zorder=5,
                   transform=ax.transAxes)
        ax.text(xf, 0.76, name, transform=ax.transAxes,
                ha="center", fontsize=15, color=col, fontweight="bold")
        ax.text(xf, 0.08, desc, transform=ax.transAxes,
                ha="center", fontsize=11, color=_M)

    ax.set_xlim(0, 1); ax.set_ylim(0, 1); ax.axis("off")
    fig.tight_layout()
    p = os.path.join(FIGS, "fig_task_familiarity.png")
    fig.savefig(p, dpi=150, bbox_inches="tight", facecolor=_D)
    plt.close(fig); return p


def fig_h2_sequential():
    """Left: weight-space diagram.  Right: Task-A accuracy vs R_A."""
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(14, 6.0))
    fig.patch.set_facecolor(_D)

    # ── left: weight space ──
    ax1.set_facecolor(_D); ax1.set_aspect("equal"); ax1.axis("off")

    θ_pre = np.array([-0.5, 0.2])
    θ_A   = np.array([1.8,  0.2])
    sh_A  = 1.4

    # Certified ball — drawn first so arrows/labels render on top
    circ = plt.Circle(θ_A, sh_A, fill=True,
                      facecolor=_H2+"22", ec=_H2, lw=2.5, ls="--")
    ax1.add_patch(circ)

    # Ball label — two short lines to keep it narrow, stays in upper-left zone
    ax1.text(θ_A[0] - sh_A + 0.10, θ_A[1] + sh_A + 0.22,
             "Certified ball\n(radius  σ½_A)",
             ha="left", fontsize=12, color=_H2)

    # θ_pre: label ABOVE the dot
    ax1.scatter(*θ_pre, s=180, color=_M, zorder=5)
    ax1.text(θ_pre[0], θ_pre[1] + 0.32, "θ_pre", fontsize=13, color=_M, ha="center")

    # θ_A: label ABOVE-RIGHT of the star
    ax1.scatter(*θ_A, s=280, color=_H2, zorder=6, marker="*")
    ax1.text(θ_A[0] + 0.16, θ_A[1] + 0.30, "θ_A", fontsize=16, color=_H2,
             fontweight="bold")

    # Phase 1 arrow (horizontal, θ_pre → θ_A)
    ax1.annotate("", xy=θ_A - np.array([0.18, 0]), xytext=θ_pre + np.array([0.12, 0]),
                 arrowprops=dict(arrowstyle="->", color=_H1, lw=2.5))
    # Phase 1 label BELOW the arrow, centred between the two endpoints
    ax1.text((θ_pre[0] + θ_A[0]) / 2, θ_pre[1] - 0.42,
             "Phase 1  →  fine-tune on Task A",
             ha="center", fontsize=12, color=_H1)

    # Phase 2: ALL LRs fine-tune on Task B from θ_A in the SAME direction.
    # Only the distance traveled differs — higher lr → farther from θ_A.
    direction = np.array([0.574, 0.819])  # [0.7, 1.0] normalised, upper-right

    # Faint dotted guide line — shows the shared direction without a crowding label
    ray_end = θ_A + 2.55 * direction
    ax1.plot([θ_A[0], ray_end[0]], [θ_A[1], ray_end[1]],
             color=_M, lw=1.2, ls=":", alpha=0.35, zorder=1)

    # Three stopping points along the same ray (distance ∝ lr)
    θ_safe   = θ_A + 0.70 * direction   # R_A = 0.5, inside ball ✓
    θ_medium = θ_A + 1.40 * direction   # R_A = 1.0, on ball boundary
    θ_unsafe = θ_A + 2.20 * direction   # R_A = 1.6, outside ball ✗

    ax1.scatter(*θ_safe,   s=180, color=_H2, zorder=6, marker="^")
    ax1.scatter(*θ_medium, s=180, color=_YL, zorder=6, marker="^")
    ax1.scatter(*θ_unsafe, s=180, color=_AC, zorder=6, marker="^")

    # safe: drop a connector below the ball → label in the open space below
    ax1.plot([θ_safe[0], θ_safe[0]], [θ_safe[1] - 0.18, -0.42],
             color=_H2, ls=":", lw=1.0, alpha=0.55, zorder=1)
    ax1.text(θ_safe[0], -0.58,
             "lr = 1e-5\nR_A < 1  ✓", ha="center", fontsize=12, color=_H2)
    # medium: to the right of the dot, slightly below
    ax1.text(θ_medium[0] + 0.20, θ_medium[1] - 0.12,
             "lr = 1e-4\nR_A ≈ 1", ha="left", fontsize=12, color=_YL)
    # unsafe: upper-right — clear of ball label (upper-left) and medium label
    ax1.text(θ_unsafe[0] + 0.14, θ_unsafe[1] + 0.16,
             "lr = 2e-4\nR_A > 1  ✗", ha="left", fontsize=12, color=_AC)

    ax1.set_xlim(-1.5, 5.2); ax1.set_ylim(-1.2, 4.2)
    ax1.set_title("Weight space  ·  Phase 1 → Task A,  Phase 2 ↗ Task B (3 LRs)",
                  fontsize=13, color=_L, pad=10)

    # ── right: accuracy vs R_A ──
    # All curves show Task-A accuracy vs R_A = ‖θ_B − θ_A‖/σ½_A.
    # H2 predicts the drop at R_A = 1 — so all inflection points should be near 1.
    # Green (lr=1e-5) is truncated to R_A ≤ 0.25: that LR never achieves high R_A;
    # showing it flat out to R_A=2.5 would falsely imply immunity at high displacement.
    _ax_style(ax2)
    R     = np.linspace(0, 2.5, 300)
    R_low = np.linspace(0, 0.25, 60)   # realistic range for lr=1e-5

    def acc(R, steepness=4.0, shift=1.0):
        return 88 / (1 + np.exp(steepness*(R - shift)))

    ax2.plot(R_low, np.full_like(R_low, 87), color=_H2, lw=2.5, label="lr = 1e-5")
    ax2.plot(R,     acc(R, 5.0, 1.00),       color=_YL, lw=2.5, label="lr = 1e-4")
    ax2.plot(R,     acc(R, 3.5, 1.00),       color=_AC, lw=2.5, label="lr = 2e-4")

    ax2.axvline(1.0, color=_L, ls="--", lw=2, alpha=0.6)
    ax2.text(1.05, 36, "R_A = 1", fontsize=13, color=_L)
    ax2.fill_betweenx([30, 93], 0,   1.0, alpha=0.07, color=_H2)
    ax2.fill_betweenx([30, 93], 1.0, 2.5, alpha=0.07, color=_AC)
    ax2.text(0.30, 91, "Safe", fontsize=13, color=_H2)
    ax2.text(1.50, 91, "Forgetting", fontsize=13, color=_AC)

    ax2.set_xlabel("R_A = ‖θ_B − θ_A‖ / σ½_A", fontsize=14, color=_L)
    ax2.set_ylabel("Task A accuracy (%)", fontsize=14, color=_L)
    ax2.set_title("Predicted: drop at R_A = 1", fontsize=14, color=_L)
    ax2.set_ylim(30, 96); ax2.set_xlim(0, 2.5)
    ax2.legend(fontsize=12, framealpha=0.2, facecolor=_B2, labelcolor=_L)

    fig.tight_layout()
    p = os.path.join(FIGS, "fig_h2_sequential.png")
    fig.savefig(p, dpi=150, bbox_inches="tight", facecolor=_D)
    plt.close(fig); return p


def fig_h3_overlap():
    """Two panels: balls overlap (merge works) vs. don't overlap (merge fails)."""
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(13, 5.8))
    fig.patch.set_facecolor(_D)

    for ax, overlap, title in [
        (ax1, True,  "d < σ½_A + σ½_B  →  merge works  ✓"),
        (ax2, False, "d > σ½_A + σ½_B  →  merge fails  ✗"),
    ]:
        ax.set_facecolor(_D); ax.set_aspect("equal"); ax.axis("off")
        gap = 1.7 if overlap else 3.0
        θ_A = np.array([-gap/2, 0]); θ_B = np.array([gap/2, 0])
        sh = 1.2

        ca = plt.Circle(θ_A, sh, fill=True, fc=_H1+"30", ec=_H1, lw=2.5)
        cb = plt.Circle(θ_B, sh, fill=True, fc=_H2+"30", ec=_H2, lw=2.5)
        ax.add_patch(ca); ax.add_patch(cb)

        ax.scatter(*θ_A, s=220, color=_H1, zorder=6, marker="*")
        ax.scatter(*θ_B, s=220, color=_H2, zorder=6, marker="*")
        ax.text(θ_A[0], θ_A[1]-0.50, "θ_A", ha="center", fontsize=15, color=_H1, fontweight="bold")
        ax.text(θ_B[0], θ_B[1]-0.50, "θ_B", ha="center", fontsize=15, color=_H2, fontweight="bold")
        ax.text(θ_A[0], θ_A[1]+sh+0.22, "σ½_A", ha="center", fontsize=13, color=_H1)
        ax.text(θ_B[0], θ_B[1]+sh+0.22, "σ½_B", ha="center", fontsize=13, color=_H2)

        θ_m = (θ_A + θ_B) / 2
        mc = _GN if overlap else _AC
        ml = "θ_merge  ✓\n(inside both balls)" if overlap else "θ_merge  ✗\n(outside both balls)"
        ax.scatter(*θ_m, s=200, color=mc, zorder=7, marker="D")
        ax.text(θ_m[0], θ_m[1]+0.48, ml, ha="center", fontsize=12, color=mc)

        # distance arrow — drawn below the ball centers
        ax.annotate("", xy=θ_B - np.array([0,0.15]), xytext=θ_A + np.array([0,-0.15]),
                    arrowprops=dict(arrowstyle="<->", color=_M, lw=1.5, alpha=0.6))
        # Distance label below centers, well clear of θ_A/θ_B name labels
        ax.text(0, -0.80, "d = ‖θ_A − θ_B‖", ha="center", fontsize=12, color=_M)

        col = _GN if overlap else _AC
        cond = "d < σ½_A + σ½_B  ✓" if overlap else "d > σ½_A + σ½_B  ✗"
        ax.text(0, -1.28, cond, ha="center", fontsize=14, color=col, fontweight="bold")
        ax.set_title(title, fontsize=14, color=col, pad=6)
        ax.set_xlim(-3.2, 3.2); ax.set_ylim(-2.1, 3.0)

    fig.tight_layout()
    p = os.path.join(FIGS, "fig_h3_overlap.png")
    fig.savefig(p, dpi=150, bbox_inches="tight", facecolor=_D)
    plt.close(fig); return p


def fig_lr_spectrum():
    """Which LRs land inside vs. outside the ball, for H1."""
    fig, ax = plt.subplots(figsize=(13, 3.8))
    fig.patch.set_facecolor(_D); ax.set_facecolor(_D); ax.axis("off")

    lrs  = ["1e-5",  "5e-5",  "1e-4",  "2e-4",   "5e-4",  "1e-3"]
    cols = [_H2,     _H2,     _H2,     _YL,      _AC,     _AC]
    lbls = ["R≪1",  "R<1",   "R<1",   "R≈1",    "R>1",   "R≫1"]
    # xs: manually spaced so lr=2e-4 (boundary marker) sits clearly RIGHT of the R=1 line
    xs   = [0.08,   0.24,    0.40,    0.64,     0.78,    0.92]

    # Horizontal axis line
    ax.axhline(0.50, xmin=0.04, xmax=0.96, color=_M, lw=2.0)

    # R=1 boundary line — placed BETWEEN lr=1e-4 (x=0.40) and lr=2e-4 (x=0.64)
    bx = 0.52
    ax.axvline(bx, ymin=0.12, ymax=0.88, color=_YL, lw=2.5, ls="--")
    ax.text(bx - 0.01, 0.95, "R = 1  boundary",
            transform=ax.transAxes, fontsize=12, color=_YL, va="top", ha="right")

    for xf, lr, col, lbl in zip(xs, lrs, cols, lbls):
        ax.scatter(xf, 0.50, s=300, color=col, zorder=5,
                   transform=ax.transAxes)
        ax.text(xf, 0.78, f"lr={lr}", transform=ax.transAxes,
                ha="center", fontsize=13, color=col, fontweight="bold")
        ax.text(xf, 0.16, lbl, transform=ax.transAxes,
                ha="center", fontsize=12, color=col)

    ax.text(0.18, 0.03, "← Inside ball (safe)", transform=ax.transAxes,
            fontsize=12, color=_H2)
    ax.text(0.62, 0.03, "Outside ball (degrades) →", transform=ax.transAxes,
            fontsize=12, color=_AC)
    ax.set_xlim(0, 1); ax.set_ylim(0, 1)
    fig.tight_layout()
    p = os.path.join(FIGS, "fig_lr_spectrum.png")
    fig.savefig(p, dpi=150, bbox_inches="tight", facecolor=_D)
    plt.close(fig); return p


# ══════════════════════════════════════════════════════════════════════════════
#  PPTX HELPERS
# ══════════════════════════════════════════════════════════════════════════════

def _prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.50)
    return prs

def _blank(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = DARK
    return sl

def _rect(sl, x, y, w, h, fill=None, line_c=None):
    shp = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else:    shp.fill.background()
    shp.line.width = Pt(0) if not line_c else Pt(1.2)
    if line_c: shp.line.color.rgb = line_c
    return shp

def _tb(sl, text, x, y, w, h, size=13, bold=False, italic=False,
        color=None, align=PP_ALIGN.LEFT):
    txb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size  = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color if color else LIGHT
    return txb

def _bullets(sl, items, x, y, w, h, size=12):
    txb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        text, lvl, col = (item,0,None) if isinstance(item,str) \
                         else (item+(None,))[:3] if len(item)==2 \
                         else item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = ("    ◦  " if lvl else "→  ") + text
        r.font.size  = Pt(size - 1 if lvl else size)
        r.font.color.rgb = col if col else (MID if lvl else LIGHT)

def _hline(sl, x, y, w, color=ACCENT, t=Pt(1.5)):
    ln = sl.shapes.add_connector(1,Inches(x),Inches(y),Inches(x+w),Inches(y))
    ln.line.color.rgb = color; ln.line.width = t

def _vbar(sl, color=ACCENT):
    shp = sl.shapes.add_shape(1,Inches(0),Inches(0),Inches(0.10),Inches(7.5))
    shp.fill.solid(); shp.fill.fore_color.rgb = color; shp.line.width = Pt(0)

def _title_bar(sl, text, subtitle=None, bar=ACCENT):
    _vbar(sl, bar)
    _tb(sl, text, 0.25, 0.10, 12.85, 0.65, size=28, bold=True)
    if subtitle:
        _tb(sl, subtitle, 0.25, 0.72, 12.85, 0.38, size=15, italic=True, color=bar)
    _hline(sl, 0.25, 1.14, 12.85, color=bar)

def _img(sl, path, x, y, w, h=None):
    if not os.path.exists(path):
        print(f"  [missing] {path}"); return
    if h: sl.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    else: sl.shapes.add_picture(path, Inches(x), Inches(y), Inches(w))

def _caption(sl, text):
    _tb(sl, text, 0.25, 6.95, 12.85, 0.45, size=12, italic=True, color=MID,
        align=PP_ALIGN.CENTER)

def _card(sl, x, y, w, h, hdr_color, hdr_text, body_items, bsize=11):
    _rect(sl, x, y, w, h, fill=BG2, line_c=hdr_color)
    _rect(sl, x, y, w, 0.46, fill=hdr_color)
    _tb(sl, hdr_text, x+0.14, y+0.06, w-0.22, 0.38,
        size=12, bold=True, color=WHITE)
    _bullets(sl, body_items, x+0.14, y+0.56, w-0.22, h-0.72, size=bsize)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDES
# ══════════════════════════════════════════════════════════════════════════════

# ── 1: Title ──────────────────────────────────────────────────────────────────
def s01_title(prs):
    sl = _blank(prs); _vbar(sl, ACCENT)
    _tb(sl, "NTRS", 0.38, 0.55, 12.5, 1.15, size=72, bold=True, color=WHITE)
    _tb(sl, "Neural Thickets  ×  Randomized Smoothing",
        0.38, 1.75, 12.0, 0.60, size=24, color=MID)
    _tb(sl, "Certified Geometry of the LLM Post-Training Weight Space",
        0.38, 2.38, 12.0, 0.50, size=18, italic=True, color=ACCENT)
    _hline(sl, 0.38, 3.05, 5.5, color=ACCENT)
    for i,(lbl,col) in enumerate([
        ("H1   Universal Adaptation Scale", H1C),
        ("H2   Forgetting Budget",          H2C),
        ("H3   Model Merging",              H3C),
    ]):
        _tb(sl, lbl, 0.38, 3.30+i*0.60, 9.0, 0.52, size=18, color=col)
    _tb(sl, "Kina Kim  |  SFU  |  June 2026",
        0.38, 6.80, 9.0, 0.45, size=14, color=MID)


# ── 2: The Problem ────────────────────────────────────────────────────────────
def s02_problem(prs):
    sl = _blank(prs); _vbar(sl, ACCENT)
    _tb(sl, "Post-training has no formal guarantees",
        0.35, 0.85, 12.5, 0.80, size=32, bold=True)
    _hline(sl, 0.35, 1.72, 12.6, color=ACCENT)

    _bullets(sl, [
        "LoRA, RLHF, and SFT adapt LLMs — but none answer: how far can we move θ before quality degrades?",
        "Practitioners tune learning rate by trial and error; there is no principled safe budget",
        "Sequential fine-tuning causes catastrophic forgetting — but when, and by how much?",
        "Model merging works empirically but has no geometric sufficient condition for success",
    ], 0.35, 1.90, 12.6, 3.20, size=16)

    _rect(sl, 0.35, 5.30, 12.6, 1.50, fill=BG2, line_c=ACCENT)
    _tb(sl, "What a certificate would buy", 0.55, 5.36, 8.0, 0.40,
        size=14, bold=True, color=ACCENT)
    _bullets(sl, [
        "Predict a safe LR and step count before running any fine-tuning",
        "Know when sequential FT will cause forgetting — without running it",
        "Certify that a model merge will succeed — from geometry alone",
    ], 0.55, 5.80, 12.2, 0.92, size=13)


# ── 3: The Bridge ─────────────────────────────────────────────────────────────
def s03_bridge(prs):
    sl = _blank(prs); _vbar(sl, ACCENT)
    _tb(sl, "We bring Randomized Smoothing into weight space",
        0.35, 0.70, 12.5, 0.75, size=30, bold=True)
    _hline(sl, 0.35, 1.50, 12.6, color=ACCENT)

    # left: Cohen vs NTRS
    _rect(sl, 0.35, 1.65, 6.0, 2.30, fill=BG2, line_c=MID)
    _tb(sl, "Cohen et al. (2019) — original RS", 0.52, 1.70, 5.6, 0.40,
        size=12, bold=True, color=MID)
    _bullets(sl, [
        "Perturb the INPUT x with Gaussian noise",
        "Certify robustness to adversarial input changes",
        "Radius in input space: r* = σ · Φ⁻¹(p_A)",
    ], 0.52, 2.15, 5.6, 1.65, size=12)

    _rect(sl, 6.65, 1.65, 6.3, 2.30, fill=BG2, line_c=H1C)
    _tb(sl, "NTRS — weight-space RS  (this work)", 6.82, 1.70, 5.9, 0.40,
        size=12, bold=True, color=H1C)
    _bullets(sl, [
        "Perturb the WEIGHTS θ with Gaussian noise",
        "Certify robustness to post-training weight changes",
        "Same Binomial CI machinery — new domain, new meaning",
    ], 6.82, 2.15, 5.9, 1.65, size=12)

    # big formula
    _rect(sl, 0.35, 4.12, 12.6, 1.40, fill=H1L, line_c=H1C)
    _tb(sl, "C(θ, σ)  =  P_{ε ~ N(0, σ²I)}[ s(θ + ε)  ≥  m ]",
        0.55, 4.26, 12.2, 0.65, size=22, bold=True, color=WHITE)
    _tb(sl, "s(θ) = quality score  (e.g. neg. WikiText-2 NLL)     m = quality threshold",
        0.55, 4.88, 12.2, 0.38, size=12, italic=True, color=MID)

    _bullets(sl, [
        "C → 1.0: model survives ALL perturbations of scale σ  (deep inside the basin)",
        "C → 0.0: model is at the edge — small noise degrades quality",
    ], 0.35, 5.68, 12.6, 0.90, size=14)


# ── 4: Density Curve (diagram) ────────────────────────────────────────────────
def s04_density_curve(prs, fig_path):
    sl = _blank(prs)
    _title_bar(sl, "C(θ, σ): How it looks on a real model", bar=H1C)
    _img(sl, fig_path, 0.30, 1.24, 8.60, 5.10)

    _rect(sl, 9.10, 1.24, 4.00, 2.60, fill=BG2, line_c=H1C)
    _tb(sl, "What the curves show", 9.28, 1.30, 3.70, 0.38, size=13, bold=True, color=H1C)
    _bullets(sl, [
        ("Pretrained (blue): peaks near 0.90, drops sharply — σ½ ≈ 0.001", 0, H1C),
        ("Random-init (red): peaks at only 0.58, decreases very slowly — σ½ ≈ 0.020", 0, ACCENT),
        "Random-init σ½ is ~20× larger — not because it is more robust, but because its loss landscape is flat and undifferentiated",
        "σ½ measures minimum SHARPNESS, not quality",
    ], 9.28, 1.74, 3.72, 2.00, size=11)

    _rect(sl, 9.10, 3.98, 4.00, 2.66, fill=BG2, line_c=YELLOW)
    _tb(sl, "The key insight", 9.28, 4.04, 3.70, 0.38, size=13, bold=True, color=YELLOW)
    _tb(sl,
        "Smaller σ½  =  sharper, better-organised minimum  =  better-trained model",
        9.28, 4.46, 3.70, 0.72, size=13, italic=True, color=LIGHT)
    _bullets(sl, [
        "Pretrained GPT-2 is HARDER to perturb away from quality than random-init",
        "Computed in ~2 min on GPU with N = 200 perturbations — no labels needed",
    ], 9.28, 5.22, 3.72, 1.30, size=11)


# ── 5: Loss basin diagram ─────────────────────────────────────────────────────
def s05_loss_basin(prs, fig_path):
    sl = _blank(prs)
    _title_bar(sl, "σ½ is the radius of the safe zone in weight space", bar=H1C)
    _img(sl, fig_path, 0.30, 1.24, 6.20, 5.60)

    _rect(sl, 6.70, 1.24, 6.30, 2.56, fill=BG2, line_c=H2C)
    _tb(sl, "Inside the ball  ✓", 6.88, 1.30, 5.90, 0.40, size=13, bold=True, color=H2C)
    _bullets(sl, [
        "‖ε‖ < σ½  →  s(θ+ε) ≥ m  with high probability",
        "Fine-tuned model stays near pretrained quality",
        "WikiText-2 PPL ratio stays ≈ 1.0",
    ], 6.88, 1.76, 5.90, 1.90, size=12)

    _rect(sl, 6.70, 3.96, 6.30, 2.84, fill=BG2, line_c=ACCENT)
    _tb(sl, "Outside the ball  ✗", 6.88, 4.02, 5.90, 0.40, size=13, bold=True, color=ACCENT)
    _bullets(sl, [
        "‖ε‖ > σ½  →  quality degrades in expectation",
        "Fine-tuned model drifts off the pretrained loss surface",
        "WikiText-2 PPL ratio rises above 1.02",
    ], 6.88, 4.48, 5.90, 1.90, size=12)

    _caption(sl, "Perturbations are isotropic Gaussian — random direction, fixed scale σ")


# ── 6: σ½ four properties ─────────────────────────────────────────────────────
def s06_sigma_props(prs):
    sl = _blank(prs)
    _title_bar(sl, "Why σ½ is the right unit for all three hypotheses", bar=H1C)

    props = [
        (H1C,  "Model-agnostic",
         ["σ½ absorbs differences in scale, arch, and pretraining quality",
          "‖Δθ‖ / σ½ is dimensionless — comparable across GPT-2, Llama, etc.",
          "Without σ½, each model needs its own threshold"]),
        (H2C,  "Measurable before fine-tuning",
         ["Computed on θ_pre in ~2 min  (no task labels, no fine-tuning)",
          "Gives a predictive budget BEFORE the experiment starts",
          "σ½_pre is reused across H1, H2, and H3 without re-running"]),
        (H3C,  "Interpretable geometry",
         ["R_A > 1 means 'you have left the certified basin of Task A'",
          "d < σ½_A + σ½_B means 'the two basins overlap — merge is safe'",
          "Geometric picture is intuitive and falsifiable"]),
        (YELLOW, "Connects to Neural Thickets",
         ["Larger / better-trained models → smaller σ½  (tighter basin)",
          "Consistent with NT solution-density scaling law",
          "σ½ is the empirical correlate of 'thicket density'"]),
    ]
    for i,(col, hdr, items) in enumerate(props):
        x = 0.25 + (i%2)*6.55
        y = 1.28 + (i//2)*3.00
        _card(sl, x, y, 6.28, 2.78, col, hdr, items, bsize=12)


# ── 7: H1 Claim ───────────────────────────────────────────────────────────────
def s07_h1_claim(prs):
    sl = _blank(prs); _vbar(sl, H1C)
    _tb(sl, "H1  ·  Universal Adaptation Scale", 0.25, 0.08, 12.5, 0.52,
        size=13, color=H1C)
    _tb(sl, "One normalised number predicts stability\nacross any model, task, or LR",
        0.25, 0.56, 12.5, 0.90, size=28, bold=True)
    _hline(sl, 0.25, 1.50, 12.85, color=H1C)

    _rect(sl, 0.25, 1.62, 12.85, 1.38, fill=H1L, line_c=H1C)
    _tb(sl, "R = ‖Δθ‖ / σ½_pre  ≈  1   is a universal stability boundary",
        0.45, 1.76, 12.4, 0.58, size=22, bold=True, color=WHITE)
    _tb(sl, "Normalising by σ½ collapses all PPL-degradation curves onto one trajectory",
        0.45, 2.28, 12.4, 0.40, size=13, italic=True, color=MID)

    _bullets(sl, [
        "R < 1  →  update inside the certified ball  →  pretraining quality preserved",
        "R > 1  →  update outside the ball  →  WikiText-2 PPL degrades",
        "σ½ is the natural unit that makes fine-tuning scale comparable across architectures",
    ], 0.25, 3.18, 12.85, 1.40, size=15)

    # Analogy row
    _tb(sl, "Analogies to known scaling laws", 0.25, 4.72, 12.85, 0.40,
        size=14, bold=True)
    analogies = [
        ("Chinchilla", "N_tokens / N_params ≈ 20  collapses loss across model sizes", H2C),
        ("μP",         "lr / √width  collapses training curves across widths",          H2C),
        ("H1  (claim)","‖Δθ‖ / σ½  collapses PPL curves across tasks, models, LRs",    H1C),
    ]
    for i,(tag, desc, col) in enumerate(analogies):
        x = 0.25 + i*4.35
        _rect(sl, x, 5.18, 4.15, 1.08, fill=BG2, line_c=col)
        _tb(sl, tag,  x+0.12, 5.24, 3.90, 0.36, size=12, bold=True, color=col)
        _tb(sl, desc, x+0.12, 5.62, 3.90, 0.55, size=11, color=LIGHT)


# ── 8: H1 Normalization diagram ───────────────────────────────────────────────
def s08_h1_norm(prs, fig_path):
    sl = _blank(prs)
    _title_bar(sl, "H1 — Visualised: before and after σ½ normalisation", bar=H1C)
    _img(sl, fig_path, 0.28, 1.24, 12.80, 5.50)
    _caption(sl, "Illustrative.  If H1 holds, the right panel collapses to a single curve with a kink at R = 1.")


# ── 9: H1 Design — scope ──────────────────────────────────────────────────────
def s09_h1_design(prs):
    sl = _blank(prs)
    _title_bar(sl, "H1  ·  Experiment Scope", bar=H1C,
               subtitle="5 models  ×  5 tasks  ×  6 LRs  =  150 fine-tuning runs")

    headers = ["Model", "Params", "Architecture", "σ½_pre"]
    rows = [
        ("GPT-2",        "124 M",  "GPT-2 decoder",    "~0.00060"),
        ("GPT-2-medium", "354 M",  "GPT-2 decoder",    "~0.00081"),
        ("GPT-2-large",  "774 M",  "GPT-2 decoder",    "~0.00060"),
        ("GPT-2-xl",     "1.5 B",  "GPT-2 decoder",    "~0.00067"),
        ("Llama-3.2-3B", "3 B",    "Llama decoder",    "~0.00020"),
    ]
    col_xs = [0.30, 3.30, 5.10, 9.60]; col_ws = [2.90, 1.70, 4.40, 2.50]

    _rect(sl, 0.25, 1.28, 12.40, 0.46, fill=H1C)
    for hdr, cx, cw in zip(headers, col_xs, col_ws):
        _tb(sl, hdr, cx, 1.30, cw, 0.40, size=11, bold=True, color=WHITE)

    for i,row in enumerate(rows):
        yr = 1.76 + i*0.48
        bg = BG2 if i%2 else RGBColor(0x12,0x22,0x38)
        _rect(sl, 0.25, yr, 12.40, 0.46, fill=bg)
        for val,cx,cw in zip(row, col_xs, col_ws):
            _tb(sl, val, cx, yr+0.05, cw, 0.38, size=12)

    # Tasks
    _tb(sl, "Tasks", 0.25, 4.26, 12.40, 0.40, size=14, bold=True)
    for i,(task,desc,col) in enumerate([
        ("SST-2",   "Binary sentiment  —  novel format",           ACCENT),
        ("MNLI",    "Textual inference  —  moderate familiarity",   YELLOW),
        ("Yahoo",   "Q&A classification  —  partial overlap",       YELLOW),
        ("DBPedia", "Entity classification  —  encyclopedic",       H1C),
        ("AG News", "News classification  —  WebText familiar",     H2C),
    ]):
        x = 0.25 + i*2.55
        _rect(sl, x, 4.72, 2.40, 0.88, fill=BG2, line_c=col)
        _tb(sl, task, x+0.10, 4.78, 2.18, 0.34, size=12, bold=True, color=col)
        _tb(sl, desc, x+0.10, 5.10, 2.18, 0.44, size=9,  color=MID)

    _bullets(sl, [
        "LoRA rank = 8 (primary)  ·  500 steps per condition  ·  eval on WikiText-2 NLL",
    ], 0.25, 5.78, 12.40, 0.40, size=13)


# ── 10: H1 Task familiarity diagram ───────────────────────────────────────────
def s10_h1_tasks(prs, fig_path):
    sl = _blank(prs)
    _title_bar(sl, "H1  ·  Why These Tasks? Controlling for Pretraining Familiarity", bar=H1C)
    _img(sl, fig_path, 0.28, 1.24, 12.80, 2.90)

    _rect(sl, 0.28, 4.28, 6.10, 2.98, fill=BG2, line_c=ACCENT)
    _tb(sl, "Familiar tasks (AG News, DBPedia)", 0.46, 4.34, 5.70, 0.40,
        size=13, bold=True, color=H2C)
    _bullets(sl, [
        "GPT-2 was pretrained on WebText — news and encyclopedic text is heavily represented",
        "Model needs only tiny parameter movement to classify these → R stays small even at high LR",
        "Risk: if we only used familiar tasks, R≈1 might hold trivially and mean nothing",
    ], 0.46, 4.80, 5.70, 2.36, size=12)

    _rect(sl, 6.58, 4.28, 6.48, 2.98, fill=BG2, line_c=ACCENT)
    _tb(sl, "Novel tasks (SST-2, MNLI)", 6.76, 4.34, 6.10, 0.40,
        size=13, bold=True, color=ACCENT)
    _bullets(sl, [
        "Short-form sentiment and textual inference are NOT dominant patterns in WebText",
        "Model must learn a genuinely new signal → large parameter movement required → R grows",
        "Including these tests whether R≈1 holds even when the model has to move far",
    ], 6.76, 4.80, 6.10, 2.36, size=12)

    _caption(sl, "Both extremes are required — without novel tasks, universality claim is untestable")


# ── 11: H1 LR & Model rationale ───────────────────────────────────────────────
def s11_h1_lr_models(prs, fig_path):
    sl = _blank(prs)
    _title_bar(sl, "H1  ·  Why These LRs and Models?", bar=H1C)
    _img(sl, fig_path, 0.28, 1.24, 12.80, 2.56)

    _rect(sl, 0.28, 4.00, 6.10, 3.24, fill=BG2, line_c=H1C)
    _tb(sl, "Learning rates: span both sides of R = 1", 0.46, 4.06, 5.70, 0.40,
        size=13, bold=True, color=H1C)
    _bullets(sl, [
        "lr = 1e-5 → R ≪ 1: the safe control condition (model barely moves)",
        "lr = 1e-4 / 2e-4 → R ≈ 0.4–0.9: expected neighbourhood of the boundary",
        "lr = 5e-4 / 1e-3 → R > 1: deliberate overshoot — confirms degradation side",
        "Without overshooting the boundary we cannot verify that σ½ predicts failure",
    ], 0.46, 4.52, 5.70, 2.60, size=12)

    _rect(sl, 6.58, 4.00, 6.48, 3.24, fill=BG2, line_c=H1C)
    _tb(sl, "Models: architecture + scale diversity", 6.76, 4.06, 6.10, 0.40,
        size=13, bold=True, color=H1C)
    _bullets(sl, [
        "GPT-2 family (124M → 1.5B): same architecture, clean parameter scaling — isolates size effect",
        "Llama-3.2-3B: different architecture AND training recipe (RLHF, instruction tuning)",
        "Without Llama, 'universal' only means 'works for GPT-2 decoder-only stack'",
        "GPT-2-xl + Llama-3B bracket the Neural Thickets 'thicket threshold' (~1.5B)",
    ], 6.76, 4.52, 6.10, 2.60, size=12)


# ── 12: H2 Claim ──────────────────────────────────────────────────────────────
def s12_h2_claim(prs):
    sl = _blank(prs); _vbar(sl, H2C)
    _tb(sl, "H2  ·  Forgetting Budget", 0.25, 0.08, 12.5, 0.52,
        size=13, color=H2C)
    _tb(sl, "The certified ball of θ_A is\na forgetting budget for Task B training",
        0.25, 0.56, 12.5, 0.90, size=26, bold=True)
    _hline(sl, 0.25, 1.52, 12.85, color=H2C)

    _rect(sl, 0.25, 1.64, 12.85, 1.34, fill=H2L, line_c=H2C)
    _tb(sl, "R_A = ‖θ_B − θ_A‖ / σ½_A  >  1   predicts Task A forgetting",
        0.45, 1.78, 12.4, 0.55, size=21, bold=True, color=WHITE)
    _tb(sl, "Whenever Phase 2 weights exit the certified ball around θ_A, Task A accuracy drops",
        0.45, 2.28, 12.4, 0.38, size=13, italic=True, color=MID)

    _bullets(sl, [
        "Continual learning: R_A = 1 is a hard budget — train on Task B until this, then stop",
        "RLHF analogy: σ½_SFT is the weight-space KL penalty radius — bounds safe RLHF movement",
        "Practical: compute σ½_A once, then predict forgetting without running Phase 2",
    ], 0.25, 3.18, 12.85, 1.60, size=15)

    _rect(sl, 0.25, 4.92, 12.85, 1.70, fill=BG2, line_c=H2C)
    _tb(sl, "Stronger claim", 0.45, 4.98, 4.0, 0.38, size=13, bold=True, color=H2C)
    _bullets(sl, [
        "Forgetting curve F(t) aligns across model sizes, task pairs, and LRs when plotted vs R_A(t)",
        "The transition at R_A = 1 is sharp — a step function, not a gradual drift",
    ], 0.45, 5.40, 12.4, 1.12, size=13)


# ── 13: H2 Sequential diagram ─────────────────────────────────────────────────
def s13_h2_diagram(prs, fig_path):
    sl = _blank(prs)
    _title_bar(sl, "H2 — Visualised: weight-space trajectory and Task A accuracy", bar=H2C)
    _img(sl, fig_path, 0.28, 1.24, 12.80, 5.54)
    _caption(sl,
        "Left: Phase 2 trajectory in weight space — low LR stays inside σ½_A ball, high LR exits. "
        "Right: predicted Task A accuracy vs R_A — drop at R_A = 1.")


# ── 14: H2 Task pair design ───────────────────────────────────────────────────
def s14_h2_tasks(prs):
    sl = _blank(prs)
    _title_bar(sl, "H2  ·  Task Pair Design — Two Requirements in Tension", bar=H2C)

    _rect(sl, 0.25, 1.28, 6.10, 5.50, fill=BG2, line_c=H2C)
    _rect(sl, 0.25, 1.28, 6.10, 0.48, fill=H2C)
    _tb(sl, "Task A  —  must converge cleanly", 0.42, 1.32, 5.80, 0.40,
        size=13, bold=True, color=WHITE)
    _bullets(sl, [
        "If Task A never converges, σ½_A reflects a pretrained (not fine-tuned) basin — the H2 formula becomes meaningless",
        "Choose a familiar task: AG News for GPT-2 (WebText = news text) → model reaches >80 % accuracy at lr = 1e-4 in 2000 steps",
        "Convergence guard: raise an error if Phase 1 accuracy < 70 %",
        "σ½_A is measured on the FINE-TUNED θ_A (not pretrained) — it reflects the Task A optimum basin",
    ], 0.42, 1.86, 5.75, 4.80, size=12)

    _rect(sl, 6.60, 1.28, 6.48, 5.50, fill=BG2, line_c=H2C)
    _rect(sl, 6.60, 1.28, 6.48, 0.48, fill=H2C)
    _tb(sl, "Task B  —  must require real parameter movement", 6.78, 1.32, 6.10, 0.40,
        size=13, bold=True, color=WHITE)
    _bullets(sl, [
        "If Task B is also familiar (e.g. another news dataset), R_A stays tiny throughout Phase 2 → no forgetting → H2 is untestable",
        "Choose a novel task: SST-2 for GPT-2 (short-form sentiment ≠ WebText distribution) → model must move θ significantly to learn it",
        "This ensures R_A grows across the LR range 1e-5 → 2e-4, giving data both inside and outside the ball",
        "Reverse for Llama-3B (SST-2 → AG News): tests whether task direction matters",
    ], 6.78, 1.86, 6.18, 4.80, size=12)


# ── 15: H2 LR/step rationale ──────────────────────────────────────────────────
def s15_h2_lr(prs):
    sl = _blank(prs)
    _title_bar(sl, "H2  ·  Why These LRs, Steps, and Ranks?", bar=H2C)

    for i,(col, hdr, items) in enumerate([
        (H2C, "Phase 1  —  lr_a = 1e-4, 2000 steps",
         ["lr = 1e-5 does NOT converge on AG News in 2000 steps — σ½_A stays pretrained-like and inflates R_A to nonsense values",
          "lr = 1e-4 produces ~4× basin widening on GPT-2 (verified in lora_density_experiment) — gives a well-characterised σ½_A",
          "2000 steps: enough for >80 % accuracy on AG News with rank = 8  (confirmed empirically)"]),
        (H2C, "Phase 2  —  three LRs, each targeting a different R_A regime",
         ["lr = 1e-5  →  R_A stays < 1 throughout 2000 steps  →  control: Task A should be preserved",
          "lr = 1e-4  →  R_A crosses 1 around step 800–1000  →  tests the transition; sharpest signal for H2",
          "lr = 2e-4  →  R_A > 1 early  →  confirms the failure side; large forgetting expected"]),
        (YELLOW, "Ranks 8 and 32 — why both?",
         ["rank = 32 moves faster in norm for the same LR (more directions updated per step)",
          "So at rank 32, R_A > 1 is reached at a lower LR — stress-tests whether rank shifts the boundary",
          "If the boundary moves with rank: the ball is not isotropic in LoRA-update space — important nuance"]),
    ]):
        _card(sl, 0.25, 1.28 + i*2.02, 12.85, 1.88, col, hdr, items, bsize=12)

    _caption(sl, "50-step eval cadence during Phase 2: fine enough to catch the R_A = 1 crossing without excessive compute")


# ── 16: H3 Claim + diagram ────────────────────────────────────────────────────
def s16_h3_claim(prs, fig_path):
    sl = _blank(prs); _vbar(sl, H3C)
    _tb(sl, "H3  ·  Model Merging", 0.25, 0.08, 12.5, 0.52,
        size=13, color=H3C)
    _tb(sl, "If the certified balls overlap, the midpoint is safe",
        0.25, 0.54, 12.5, 0.70, size=28, bold=True)
    _hline(sl, 0.25, 1.30, 12.85, color=H3C)

    _rect(sl, 0.25, 1.42, 12.85, 0.96, fill=H3L, line_c=H3C)
    _tb(sl, "d(θ_A, θ_B) < σ½_A + σ½_B   →   merging θ_A and θ_B preserves both tasks",
        0.45, 1.55, 12.4, 0.50, size=20, bold=True, color=WHITE)
    _tb(sl, "Geometric sufficient condition: if both models lie inside each other's certified ball, their midpoint does too",
        0.45, 2.04, 12.4, 0.28, size=12, italic=True, color=MID)

    _img(sl, fig_path, 0.30, 2.48, 8.20, 4.70)

    _rect(sl, 8.68, 2.48, 4.40, 4.70, fill=BG2, line_c=H3C)
    _tb(sl, "Why this is novel", 8.86, 2.54, 4.10, 0.38, size=13, bold=True, color=H3C)
    _bullets(sl, [
        "All prior merge methods (TIES, DARE, task arithmetic) are empirical — no geometric certificate",
        "H3 predicts success BEFORE running the merge — from σ½_A, σ½_B, and ‖θ_A − θ_B‖ alone",
        "If H1 holds, keeping ‖Δθ_A‖ < σ½ during training guarantees d < σ½_A + σ½_B",
        "The certificate is conservative: gap > 0 is sufficient, not necessary",
    ], 8.86, 2.96, 4.12, 4.10, size=11)


# ── 17: H3 Task & model choices ───────────────────────────────────────────────
def s17_h3_design(prs):
    sl = _blank(prs)
    _title_bar(sl, "H3  ·  Why These Tasks, Model, and Merge Method?", bar=H3C)

    _card(sl, 0.25, 1.28, 4.15, 5.50, H3C, "Why GPT-2-medium?",
          ["GPT-2 (124M) is capacity-limited — multi-task performance saturates at low accuracy, masking the merge signal",
           "GPT-2-medium (354M) has enough capacity to hold two tasks simultaneously — we can actually measure success",
           "σ½_pre ≈ 0.00081 already characterised from H1; baseline well-established",
           "Cheap enough to run 9 (LR_A, LR_B) conditions + merge evals in one A40 session"], bsize=12)

    _card(sl, 4.58, 1.28, 4.15, 5.50, H3C, "Why AG News + SST-2?",
          ["Different label spaces (4-class vs binary) + different text formats → genuinely conflicting weight updates",
           "AG News is familiar → θ_A stays close to θ_pre → σ½_A stays large → easier for condition to be satisfied",
           "SST-2 is novel → θ_B must move further → harder to satisfy condition → tests the boundary",
           "If both tasks were familiar, d would always be small → certificate always satisfied → H3 untestable",
           "A 'boring' pair never produces failures — we NEED gap < 0 cells to falsify H3"], bsize=12)

    _card(sl, 8.91, 1.28, 4.17, 5.50, YELLOW, "Why simple averaging?",
          ["θ_merge = (θ_A + θ_B) / 2 is the midpoint in weight space — exactly what the ball-overlap condition predicts",
           "If the condition works for the simplest merge method, it is the strongest claim",
           "Exotic methods (TIES, DARE) re-weight parameters in ways that obscure the geometry",
           "If simple averaging FAILS where gap > 0, H3 is cleanly falsified",
           "If simple averaging SUCCEEDS where gap > 0, H3 is confirmed by the most skeptical test"], bsize=12)


# ── 18: H3 LR Grid ────────────────────────────────────────────────────────────
def s18_h3_grid(prs):
    sl = _blank(prs)
    _title_bar(sl, "H3  ·  The 3×3 LR Grid — Tracing the Merge Boundary", bar=H3C,
               subtitle="We need cells where the certificate is satisfied AND violated — otherwise we cannot falsify H3")

    # Grid
    lr_labels = ["lr_B = 1e-5", "lr_B = 1e-4", "lr_B = 2e-4"]
    row_labels = ["lr_A = 1e-5", "lr_A = 1e-4", "lr_A = 2e-4"]
    cells = [
        # (text, fill_color, border_color)
        ("d ≪ σ½_A + σ½_B\ngap >> 0\nH3 predicts ✓ SUCCESS\n(control cell)",    H2L, H2C),
        ("d moderate\ngap > 0\nH3 predicts ✓ SUCCESS",                            H2L, H2C),
        ("d grows\ngap ≈ 0\nnear-boundary case",                                  BG2, YELLOW),
        ("d moderate\ngap > 0\nH3 predicts ✓ SUCCESS",                            H2L, H2C),
        ("d large\ngap ≈ 0\nmost interesting cell",                                BG2, YELLOW),
        ("d very large\ngap < 0\nH3 predicts ✗ FAILURE",                          RGBColor(0x3a,0x0c,0x0c), ACCENT),
        ("d grows\ngap ≈ 0\nasymmetric boundary",                                  BG2, YELLOW),
        ("d large\ngap < 0\nH3 predicts ✗ FAILURE",                               RGBColor(0x3a,0x0c,0x0c), ACCENT),
        ("d ≫ σ½_A + σ½_B\ngap ≪ 0\nH3 predicts ✗ FAILURE\n(stress-test cell)", RGBColor(0x3a,0x0c,0x0c), ACCENT),
    ]

    col_w = 3.60; row_h = 1.62
    col_xs = [1.65, 5.35, 9.05]
    row_ys = [1.74, 3.46, 5.18]

    # Column headers
    _rect(sl, 1.55, 1.28, 11.2, 0.42, fill=H3L)
    for j, lbl in enumerate(lr_labels):
        _tb(sl, lbl, col_xs[j], 1.30, col_w, 0.38, size=12, bold=True,
            color=H3C, align=PP_ALIGN.CENTER)

    # Row labels
    for i, lbl in enumerate(row_labels):
        _tb(sl, lbl, 0.10, row_ys[i]+0.50, 1.42, 0.56, size=11, bold=True,
            color=H3C, align=PP_ALIGN.CENTER)

    # Cells
    for k, (text, fc, bc) in enumerate(cells):
        i, j = divmod(k, 3)
        x = col_xs[j]; y = row_ys[i]
        _rect(sl, x, y, col_w-0.10, row_h-0.08, fill=fc, line_c=bc)
        _tb(sl, text, x+0.14, y+0.12, col_w-0.30, row_h-0.20,
            size=10, color=LIGHT)

    _caption(sl,
        "Green = H3 predicts success  ·  Yellow = near boundary  ·  Red = H3 predicts failure  "
        "·  Key test: does the merge success map align with the gap sign?")


# ── 19: Summary ───────────────────────────────────────────────────────────────
def s19_summary(prs):
    sl = _blank(prs); _vbar(sl, ACCENT)
    _tb(sl, "Three Geometric Bets on σ½", 0.25, 0.12, 12.5, 0.65,
        size=30, bold=True, color=WHITE)
    _hline(sl, 0.25, 0.82, 12.85, color=ACCENT)

    rows = [
        (H1C, "H1  Adaptation Scale",
         "R = ‖Δθ‖/σ½ ≈ 1 is a universal stability boundary",
         "If confirmed →", "A model-agnostic LR / step budget rule, pre-hoc, no fine-tuning needed"),
        (H2C, "H2  Forgetting Budget",
         "R_A > 1 predicts Task A forgetting after sequential FT",
         "If confirmed →", "σ½_A is a certified continual-learning budget; direct weight-space analog of KL penalty"),
        (H3C, "H3  Model Merging",
         "d(θ_A, θ_B) < σ½_A + σ½_B predicts merge success",
         "If confirmed →", "First geometry-based certificate for safe model merging — predict before running"),
    ]
    for i,(col, title, claim, tag, impact) in enumerate(rows):
        y = 1.00 + i*2.10
        _rect(sl, 0.25, y, 12.85, 1.95, fill=BG2)
        _rect(sl, 0.25, y, 0.28, 1.95, fill=col)
        _tb(sl, title, 0.68, y+0.10, 3.60, 0.45, size=15, bold=True, color=col)
        _tb(sl, claim, 0.68, y+0.58, 5.80, 0.46, size=14, color=LIGHT)
        _rect(sl, 6.75, y+0.28, 6.10, 1.30, fill=RGBColor(0x10,0x10,0x20), line_c=col)
        _tb(sl, tag,    6.92, y+0.34, 5.80, 0.32, size=11, bold=True, color=col)
        _tb(sl, impact, 6.92, y+0.68, 5.80, 0.80, size=11, italic=True, color=MID)

    _rect(sl, 0.25, 7.20, 12.85, 0.20, fill=BG2)
    _tb(sl, "Shared infrastructure: same σ½ estimator  ·  same GPU pipeline  ·  density sweeps reused across all three hypotheses",
        0.45, 7.20, 12.4, 0.22, size=11, color=MID, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════════════════════
def main():
    print("Generating figures …")
    f_density  = fig_density_curve()
    f_basin    = fig_loss_basin()
    f_h1_norm  = fig_h1_normalization()
    f_tasks    = fig_task_familiarity()
    f_h2       = fig_h2_sequential()
    f_h3       = fig_h3_overlap()
    f_lr       = fig_lr_spectrum()
    print("  done.")

    prs = _prs()
    s01_title(prs)
    s02_problem(prs)
    s03_bridge(prs)
    s04_density_curve(prs, f_density)
    s05_loss_basin(prs, f_basin)
    s06_sigma_props(prs)
    s07_h1_claim(prs)
    s08_h1_norm(prs, f_h1_norm)
    s09_h1_design(prs)
    s10_h1_tasks(prs, f_tasks)
    s11_h1_lr_models(prs, f_lr)
    s12_h2_claim(prs)
    s13_h2_diagram(prs, f_h2)
    s14_h2_tasks(prs)
    s15_h2_lr(prs)
    s16_h3_claim(prs, f_h3)
    s17_h3_design(prs)
    s18_h3_grid(prs)
    s19_summary(prs)

    prs.save(OUT)
    print(f"Saved → {OUT}  ({len(prs.slides)} slides)")

if __name__ == "__main__":
    main()
