"""
make_presentation2.py
NTRS follow-up presentation — figure-centric rebuild.
Run make_pres_figures.py first to generate results/pres_figures/*.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ───────────────────────────────────────────────────────────────────
DARK    = RGBColor(0x1a, 0x1a, 0x2e)
BG2     = RGBColor(0x0d, 0x1b, 0x2a)
ACCENT  = RGBColor(0xe9, 0x4f, 0x37)
ACCENT2 = RGBColor(0x39, 0x7b, 0xbf)
LIGHT   = RGBColor(0xf5, 0xf5, 0xf5)
MID     = RGBColor(0xb0, 0xb8, 0xc8)
GREEN   = RGBColor(0x43, 0xc5, 0x9e)
YELLOW  = RGBColor(0xf5, 0xc5, 0x18)

W = Inches(13.33)
H = Inches(7.5)

ROOT    = os.path.dirname(os.path.abspath(__file__))
FIGS    = os.path.join(ROOT, "results", "pres_figures")

def fig(name):
    return os.path.join(FIGS, name)


# ── Helpers ───────────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color=DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def title_bar(slide, text, subtitle=None):
    """Standard title bar: accent stripe + title + subtitle + separator."""
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.1), H)
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    txbox(slide, text, Inches(0.25), Inches(0.12), Inches(12.8), Inches(0.62),
          size=32, bold=True)
    if subtitle:
        txbox(slide, subtitle, Inches(0.25), Inches(0.75), Inches(12.8), Inches(0.38),
              size=18, italic=True, color=ACCENT)
    hline(slide, Inches(0.25), Inches(1.18), Inches(12.8), color=ACCENT)

def txbox(slide, text, x, y, w, h,
          size=22, bold=False, color=LIGHT, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = color
    return tb

def hline(slide, x, y, w, color=ACCENT, thickness=Pt(2)):
    ln = slide.shapes.add_connector(1, x, y, x + w, y)
    ln.line.color.rgb = color
    ln.line.width     = thickness

def bullet_box(slide, items, x, y, w, h, size=18, color=LIGHT, bullet="→"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        p.space_after = Pt(9)

def add_figure(slide, fname, x, y, w, h):
    path = fig(fname)
    if os.path.exists(path):
        slide.shapes.add_picture(path, x, y, w, h)
    else:
        # placeholder box so we know what's missing
        sh = slide.shapes.add_shape(1, x, y, w, h)
        sh.fill.solid(); sh.fill.fore_color.rgb = BG2
        sh.line.color.rgb = MID
        txbox(slide, f"[missing: {fname}]", x + Inches(0.1), y + Inches(0.1),
              w - Inches(0.2), h - Inches(0.2), size=14, color=MID)

def caption_bar(slide, text, color=MID, size=15):
    """Bottom caption strip."""
    txbox(slide, text, Inches(0.25), Inches(6.95), Inches(12.8), Inches(0.45),
          size=size, color=color, italic=True, align=PP_ALIGN.CENTER)

def code_box(slide, lines, x, y, w, h, size=16):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = BG2
    sh.line.fill.background()
    tb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1),
                                  w - Inches(0.3), h - Inches(0.2))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run(); run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = LIGHT
        run.font.name = "Courier New"
        p.space_after = Pt(4)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
def slide_title(prs):
    s = blank_slide(prs); bg(s)
    bar = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.12), H)
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    txbox(s, "Certified Basin Geometry at Scale",
          Inches(0.4), Inches(1.1), Inches(12.5), Inches(1.3),
          size=44, bold=True)
    txbox(s, "How Fine-Tuning Method and Model Scale Determine Weight-Space Robustness",
          Inches(0.4), Inches(2.6), Inches(11), Inches(0.65),
          size=24, color=MID, italic=True)
    hline(s, Inches(0.4), Inches(3.5), Inches(6), color=ACCENT)
    txbox(s, "Neural Thickets  ×  Randomized Smoothing  |  Follow-up",
          Inches(0.4), Inches(3.75), Inches(9), Inches(0.55),
          size=21, color=ACCENT2)
    txbox(s, "Kina Kim  |  SFU  |  June 2026",
          Inches(0.4), Inches(6.5), Inches(9), Inches(0.5),
          size=17, color=MID)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Framework Recap (density curve as visual)
# ══════════════════════════════════════════════════════════════════════════════
def slide_recap(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "Quick Recap — Framework & Key Definitions",
              "What we measure and why")

    # Left: density curve figure (GPT-2 pretrained + LoRA, 2-panel top row)
    # Use fig1 top-left panel as visual anchor — or just use fig6 left panel
    # Actually fig6 is the best compact single-panel demo
    add_figure(s, "fig6_full_ft_comparison.png",
               Inches(0.25), Inches(1.3), Inches(7.8), Inches(3.3))

    # Right: framework equations
    txbox(s, "Framework", Inches(8.3), Inches(1.3), Inches(4.8), Inches(0.4),
          size=18, bold=True, color=ACCENT2)
    code_box(s, [
        "C(θ, σ) = P[NLL(θ+ε) ≤ NLL(θ) + slack]",
        "          ε ~ N(0, σ²I)",
        "",
        "σ½  = σ where C drops to peak / 2",
        "     → certified basin WIDTH",
        "",
        "norm/σ½ = ‖θ_ft − θ_pre‖ / σ½_pre",
        "         → update relative to basin",
    ], Inches(8.3), Inches(1.75), Inches(4.8), Inches(2.8), size=14)

    # Bottom: new question
    hline(s, Inches(0.25), Inches(4.75), Inches(12.8), color=MID, thickness=Pt(1))
    txbox(s, "Previous (GPT-2 only):", Inches(0.25), Inches(4.9),
          Inches(3.5), Inches(0.4), size=16, bold=True, color=ACCENT2)
    bullet_box(s, [
        "norm/σ½ < 1  →  PPL stable,  basin widens 1.1–1.75×",
        "Full FT (rank=768): barely widens (1.002×) — low-rank constraint drives effect",
    ], Inches(3.9), Inches(4.9), Inches(9.0), Inches(0.9), size=15, color=MID)

    txbox(s, "New question:", Inches(0.25), Inches(5.85),
          Inches(2.5), Inches(0.4), size=17, bold=True, color=YELLOW)
    txbox(s, "Does the certified ball boundary hold across GPT-2, Qwen2.5-3B, Llama-3.2-3B, and Llama-3.1-8B?",
          Inches(2.9), Inches(5.85), Inches(10.0), Inches(0.45),
          size=17, color=LIGHT)
    txbox(s, "Does LoRA basin widening persist at 3B and 8B parameter scale?",
          Inches(2.9), Inches(6.35), Inches(10.0), Inches(0.4),
          size=17, color=LIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Experiment Verification: Density Curves
# ══════════════════════════════════════════════════════════════════════════════
def slide_density_verification(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "Experiment Verification — Density Curves Across 4 Models",
              "C(θ, σ) before and after LoRA fine-tuning (lr=1e-4, rank=8)")

    # Full-width 4-panel density figure
    add_figure(s, "fig1_density_grid.png",
               Inches(0.25), Inches(1.28), Inches(12.83), Inches(5.52))

    caption_bar(s,
        "Each panel: black=pretrained, dashed=LoRA lr=1e-4.  "
        "Dotted vertical = σ½ (50% density point).  "
        "Note: σ½_pre shrinks 5× from GPT-2 → 8B.  "
        "Near-identical pretrained vs LoRA curves at large scale = tight ball with little room to widen.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Experiment Matrix
# ══════════════════════════════════════════════════════════════════════════════
def slide_matrix(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "Complete Experiment Matrix — All 4 Models × 6 LRs",
              "Widening ratio / update size / PPL ratio at a glance  |  Red line = certified ball boundary")

    # Main heatmap matrix (wide figure)
    add_figure(s, "fig2_lr_matrix.png",
               Inches(0.25), Inches(1.28), Inches(12.83), Inches(4.02))

    # Summary bullets below
    hline(s, Inches(0.25), Inches(5.45), Inches(12.83), color=MID, thickness=Pt(0.75))
    bullet_box(s, [
        "Left panel: σ½ widening ratio — green >1 (widening), yellow ≈1 (flat), red <1 (narrowing)",
        "Center panel: log₁₀(norm/σ½) — red line marks certified ball boundary (norm/σ½=1)",
        "Right panel: PPL ratio — strongly tracks ball position: inside ball → PPL ≈ 1.0, outside → degrades",
        "lr=1e-3: catastrophic collapse across all models (PPL +160% to +380%, norm/σ½ 4×–13×)",
    ], Inches(0.25), Inches(5.55), Inches(12.83), Inches(1.65), size=15, color=MID)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Finding 1: σ½ Scaling Law
# ══════════════════════════════════════════════════════════════════════════════
def slide_scaling_law(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "Finding 1 — σ½ Shrinks with Model Capability",
              "Better models sit in tighter weight-space minima")

    # Left: bar chart (w=9, h=5.0 → aspect 10/5.5=1.82, so h=9/1.82=4.95)
    add_figure(s, "fig3_sigma_scaling.png",
               Inches(0.25), Inches(1.28), Inches(9.0), Inches(4.95))

    # Right: annotation panel
    txbox(s, "Key numbers", Inches(9.5), Inches(1.3), Inches(3.6), Inches(0.4),
          size=17, bold=True, color=ACCENT2)
    code_box(s, [
        "GPT-2  (124M):  σ½=7.78e-4",
        "Qwen   (2.5B):  σ½=2.96e-4  (2.6×)",
        "Llama  (3.0B):  σ½=1.95e-4  (4.0×)",
        "Llama  (8.0B):  σ½=1.57e-4  (5.0×)",
    ], Inches(9.5), Inches(1.75), Inches(3.6), Inches(1.85), size=13)

    txbox(s, "Interpretation", Inches(9.5), Inches(3.8), Inches(3.6), Inches(0.4),
          size=17, bold=True, color=ACCENT2)
    bullet_box(s, [
        "σ½ tracks pretrained PPL — smaller basin = better model",
        "Consistent with Neural Thickets: denser solution clusters at larger scale",
        "Tighter σ½ → safe LR budget shrinks with scale",
        "Predicts: thicket emergence (~1.5B) should show a kink in this curve",
    ], Inches(9.5), Inches(4.25), Inches(3.6), Inches(2.4), size=14, color=MID)

    caption_bar(s,
        "σ½_pre measured on pretrained model only — no fine-tuning.  "
        "Dashed line = log-linear trend.  GPT-2 is 5× wider than Llama-3.1-8B.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Finding 2: Universal PPL Safety Boundary (THE KEY RESULT)
# ══════════════════════════════════════════════════════════════════════════════
def slide_phase_boundary(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "Finding 2 — Universal PPL Safety Boundary  (Method-Agnostic)",
              "norm/σ½ < 1  →  PPL stable     LoRA rank=8  AND  Full FT rank=768  |  4 models  |  28+ conditions")

    # Large phase boundary figure — the centerpiece
    add_figure(s, "fig4_phase_boundary.png",
               Inches(0.25), Inches(1.28), Inches(9.2), Inches(5.3))

    # Right: annotation
    txbox(s, "The rule", Inches(9.65), Inches(1.3), Inches(3.45), Inches(0.4),
          size=17, bold=True, color=GREEN)
    code_box(s, [
        "norm/σ½ < 1",
        "  → PPL ratio ≈ 1.000  ✓",
        "",
        "norm/σ½ > 1",
        "  → PPL ratio > 1.02  ✗",
    ], Inches(9.65), Inches(1.75), Inches(3.45), Inches(1.8), size=14)

    txbox(s, "Evidence", Inches(9.65), Inches(3.75), Inches(3.45), Inches(0.4),
          size=17, bold=True, color=ACCENT2)
    bullet_box(s, [
        "28+ data points across 4 models, 2 methods, zero exceptions",
        "● LoRA rank=8  (24 pts, 4 models × 6 LRs)",
        "▲ Full FT rank=768  (4 pts, GPT-2 + Llama-3B LR sweep)",
        "Boundary is σ½-normalized — absolute LR is irrelevant",
        "Method doesn't matter — only norm/σ½ matters",
    ], Inches(9.65), Inches(4.2), Inches(3.45), Inches(2.35), size=12, color=MID)

    caption_bar(s,
        "Circles = LoRA rank=8.  Triangles = Full FT rank=768.  "
        "Both methods: PPL ≈ 1 inside the ball, rises outside.  "
        "Green = safe zone.  Pink = PPL-degraded zone.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Finding 3: Scale-Dependent Widening
# ══════════════════════════════════════════════════════════════════════════════
def slide_widening(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "Finding 3 — Widening Is a Small-Model Phenomenon",
              "GPT-2 widens 1.1–1.75× inside the ball; 3B+ models show flat or narrowing")

    # Main widening plot
    add_figure(s, "fig5_widening_vs_lr.png",
               Inches(0.25), Inches(1.28), Inches(9.0), Inches(5.08))

    # Right annotation
    txbox(s, "Best safe widening", Inches(9.5), Inches(1.3), Inches(3.6), Inches(0.4),
          size=16, bold=True, color=ACCENT2)
    code_box(s, [
        "GPT-2  (124M): 1.75×",
        "  (lr=2e-4, norm/σ½=0.21)",
        "",
        "Qwen   (2.5B): 0.82×",
        "  (lr=1e-5, norm/σ½=0.06)",
        "",
        "Llama  (3.0B): 0.99×",
        "  (lr=5e-5, norm/σ½=0.12)",
        "",
        "Llama  (8.0B): 1.03×",
        "  (lr=1e-4, norm/σ½=0.58)",
    ], Inches(9.5), Inches(1.75), Inches(3.6), Inches(3.0), size=12)

    txbox(s, "Why?", Inches(9.5), Inches(4.95), Inches(3.6), Inches(0.35),
          size=16, bold=True, color=ACCENT2)
    bullet_box(s, [
        "Tighter σ½ at scale → less room to widen inside ball",
        "Llama-8B: norm/σ½=0.58 already at lr=1e-4 — little space left",
        "LoRA rank=8 may not have the capacity to reshape a tighter basin",
    ], Inches(9.5), Inches(5.35), Inches(3.6), Inches(1.9), size=12, color=MID)

    caption_bar(s,
        "Open markers = norm/σ½ > 1 (escaped certified ball).  "
        "Widening above the dashed line = basin expanded after fine-tuning.  "
        "Clipped at 5× for readability.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Full FT Control (Mechanistic Evidence)
# ══════════════════════════════════════════════════════════════════════════════
def slide_full_ft(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "Method Comparison — Same Model, Same Norm Budget, Different Method",
              "Llama-3.2-3B: LoRA rank=8 vs Full FT rank=768  |  F2 confirmed method-agnostic")

    # Two-panel method comparison figure
    add_figure(s, "fig7_method_comparison.png",
               Inches(0.25), Inches(1.28), Inches(9.8), Inches(5.0))

    # Right column: key numbers
    txbox(s, "Llama-3B results", Inches(10.2), Inches(1.3), Inches(2.9), Inches(0.4),
          size=15, bold=True, color=ACCENT2)
    code_box(s, [
        "LoRA rank=8:",
        "  lr=1e-4, norm/σ½=0.22",
        "  ratio=0.955×  (narrows)",
        "  PPL=1.007  ✓",
        "",
        "Full FT rank=768:",
        "  lr=1e-5, norm/σ½=0.37",
        "  ratio=1.184×  (widens!)",
        "  PPL=1.009  ✓",
        "",
        "Both inside ball →",
        "  both PPL-stable",
    ], Inches(10.2), Inches(1.75), Inches(2.9), Inches(3.2), size=11)

    txbox(s, "Key insight", Inches(10.2), Inches(5.1), Inches(2.9), Inches(0.35),
          size=14, bold=True, color=GREEN)
    bullet_box(s, [
        "F2 holds for both methods",
        "Widening reversal: at 3B scale, full FT widens, LoRA narrows",
        "Low-rank claim holds for GPT-2, but reverses at 3B",
    ], Inches(10.2), Inches(5.5), Inches(2.9), Inches(1.7), size=11, color=MID)

    caption_bar(s,
        "Left: density curves — full FT shifts σ½ right more than LoRA.  "
        "Right: widening scatter — full FT widens (▲) while LoRA (●) is flat.  "
        "Both at Llama-3B, SST-2, 500 steps.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — RandOpt Negative Result
# ══════════════════════════════════════════════════════════════════════════════
def slide_randopt(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "RandOpt Exploration — Negative Result",
              "Random perturbation selection fails without a meaningful task scoring signal")

    # Left: what we tried + what happened
    txbox(s, "Setup  (GPT-2, N=500, K=10)", Inches(0.25), Inches(1.3),
          Inches(5.8), Inches(0.4), size=17, bold=True, color=ACCENT2)
    bullet_box(s, [
        "Sample N=500 Gaussian weight perturbations at each σ",
        "Score each by SST-2 zero-shot accuracy",
        "Keep top-K=10 deltas, average them  →  θ_ft",
        "Measure σ½ before and after; compare to LoRA",
    ], Inches(0.25), Inches(1.75), Inches(5.8), Inches(1.8), size=15)

    txbox(s, "What happened", Inches(0.25), Inches(3.65),
          Inches(5.8), Inches(0.4), size=17, bold=True, color=ACCENT)
    code_box(s, [
        "GPT-2 SST-2 zero-shot accuracy:  0.60 (near random)",
        "Top-K mean accuracy:              0.60 (no improvement)",
        "σ½ appears to 'widen' (1.4×–2.2×) — but meaningless.",
        "We averaged random noise, not selected good solutions.",
    ], Inches(0.25), Inches(4.1), Inches(5.8), Inches(1.65), size=14)

    txbox(s, "Root cause: no scoring signal",
          Inches(0.25), Inches(5.85), Inches(5.8), Inches(0.4),
          size=15, bold=True, color=ACCENT)
    txbox(s, "GPT-2 base can't do SST-2 zero-shot reliably → all perturbations score identically → "
             "selection is random → RandOpt degenerates to noise averaging.",
          Inches(0.25), Inches(6.3), Inches(5.8), Inches(0.6),
          size=13, color=MID)

    # Right: fix + NT comparison
    hline(s, Inches(6.4), Inches(1.28), Inches(0), color=MID, thickness=Pt(0.75))
    vbar = s.shapes.add_shape(1, Inches(6.38), Inches(1.28),
                               Inches(0.04), Inches(5.85))
    vbar.fill.solid(); vbar.fill.fore_color.rgb = MID
    vbar.line.fill.background()

    txbox(s, "NT paper context", Inches(6.6), Inches(1.3), Inches(6.5), Inches(0.4),
          size=17, bold=True, color=ACCENT2)
    bullet_box(s, [
        "Gan & Isola used 7B+ models on math (GSM8K) — tasks where "
        "the base model already scores meaningfully (30%+ vs 5% random)",
        "Selection signal is what makes RandOpt work",
    ], Inches(6.6), Inches(1.75), Inches(6.5), Inches(1.3), size=15, color=MID)

    txbox(s, "Two paths forward", Inches(6.6), Inches(3.2), Inches(6.5), Inches(0.4),
          size=17, bold=True, color=GREEN)
    bullet_box(s, [
        "Option A: Use supervised scoring (cross-entropy on labeled examples) — "
        "gives RandOpt the same task signal as LoRA.  Fair comparison.",
        "Option B: Focus on Claim 1 (ball boundary as universal LR rule) as primary "
        "contribution.  Treat RandOpt as future work.",
    ], Inches(6.6), Inches(3.65), Inches(6.5), Inches(2.0), size=15)

    txbox(s, "Recommended: Option B for now — Claim 1 is the stronger, cleaner result across 4 models.",
          Inches(6.6), Inches(5.85), Inches(6.5), Inches(0.55),
          size=14, bold=True, color=YELLOW)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Summary & Open Questions
# ══════════════════════════════════════════════════════════════════════════════
def slide_summary(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "Summary & Open Questions")

    # Left column: findings
    txbox(s, "What we have", Inches(0.25), Inches(1.28),
          Inches(6.5), Inches(0.42), size=20, bold=True, color=GREEN)

    findings = [
        ("F1  σ½ Scaling Law",
         "σ½_pre shrinks 5× from GPT-2 → 8B. Better models = tighter basins. Aligns with Neural Thickets."),
        ("F2  Universal Safety Boundary  (method-agnostic)",
         "norm/σ½ < 1 predicts PPL stability. Confirmed for LoRA rank=8 (24 pts) AND full FT rank=768 (4 pts). Zero exceptions across 28+ conditions, 4 models, 2 methods."),
        ("F3  Scale + Method Dependent Widening",
         "GPT-2 LoRA: widens 1.1–1.75×.  3B LoRA: flat/narrows (0.95×).  3B full FT: widens (1.18×).  Widening depends on both scale and method."),
        ("F4  Low-Rank Effect  (GPT-2 scale observation)",
         "At GPT-2 scale: full FT rank=768 barely widens (1.002×), LoRA rank=8 widens (1.12–1.29×). Pattern reverses at 3B — full FT widens, LoRA narrows."),
    ]
    for i, (title, body) in enumerate(findings):
        y = Inches(1.75) + Inches(1.2)*i
        sh = s.shapes.add_shape(1, Inches(0.25), y, Inches(6.5), Inches(1.1))
        sh.fill.solid(); sh.fill.fore_color.rgb = BG2
        sh.line.color.rgb = ACCENT2
        txbox(s, title, Inches(0.4), y + Inches(0.06), Inches(6.2), Inches(0.38),
              size=16, bold=True, color=ACCENT2)
        txbox(s, body, Inches(0.4), y + Inches(0.44), Inches(6.2), Inches(0.58),
              size=13, color=MID)

    # Divider
    vbar = s.shapes.add_shape(1, Inches(7.08), Inches(1.28), Inches(0.04), Inches(5.85))
    vbar.fill.solid(); vbar.fill.fore_color.rgb = MID
    vbar.line.fill.background()

    # Right column: open questions
    txbox(s, "Open Questions", Inches(7.3), Inches(1.28),
          Inches(5.8), Inches(0.42), size=20, bold=True, color=ACCENT)
    bullet_box(s, [
        "Why does GPT-2 widen but 3B doesn't?  Architecture, scale, or pretraining data diversity?",
        "Does a wider σ½ after FT improve downstream generalization — or is it a geometric artifact?",
        "Thicket threshold (~1.5B, per NT paper): should σ½ show a kink there?",
        "RandOpt with supervised scoring — does it naturally stay inside the certified ball?",
        "Can σ½ serve as an overfitting diagnostic during training (per professor's suggestion)?",
    ], Inches(7.3), Inches(1.78), Inches(5.8), Inches(3.5), size=15)

    hline(s, Inches(0.25), Inches(6.42), Inches(12.83), color=ACCENT, thickness=Pt(1))
    txbox(s,
          "For professor:  Is 'certified ball as universal LR budget rule' the right paper framing?  "
          "Worth pursuing supervised RandOpt for Claim 2?",
          Inches(0.25), Inches(6.55), Inches(12.83), Inches(0.8),
          size=17, bold=True, color=YELLOW)


# ══════════════════════════════════════════════════════════════════════════════
# BUILD
# ══════════════════════════════════════════════════════════════════════════════
if __name__ == "__main__":
    prs = new_prs()
    slide_title(prs)
    slide_recap(prs)
    slide_density_verification(prs)
    slide_matrix(prs)
    slide_scaling_law(prs)
    slide_phase_boundary(prs)
    slide_widening(prs)
    slide_full_ft(prs)
    slide_randopt(prs)
    slide_summary(prs)

    out = os.path.join(ROOT, "ntrs_presentation2.pptx")
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")
