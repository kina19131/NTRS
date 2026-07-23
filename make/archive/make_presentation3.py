"""
make_presentation3.py
NTRS progress-update presentation — incorporates steps sweep, trajectory run,
full FT Llama-3B sweep, and null-space results.
Run make_pres_figures.py first to regenerate results/pres_figures/*.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ───────────────────────────────────────────────────────────────────
DARK   = RGBColor(0x1a, 0x1a, 0x2e)
BG2    = RGBColor(0x0d, 0x1b, 0x2a)
ACCENT = RGBColor(0xe9, 0x4f, 0x37)
ACC2   = RGBColor(0x39, 0x7b, 0xbf)
LIGHT  = RGBColor(0xf5, 0xf5, 0xf5)
MID    = RGBColor(0xb0, 0xb8, 0xc8)
GREEN  = RGBColor(0x43, 0xc5, 0x9e)
YELLOW = RGBColor(0xf5, 0xc5, 0x18)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)

W = Inches(13.33)
H = Inches(7.5)

ROOT = os.path.dirname(os.path.abspath(__file__))
FIGS = os.path.join(ROOT, "results", "pres_figures")


def fig(name):
    return os.path.join(FIGS, name)


# ── Helpers ───────────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color=DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def title_bar(slide, text, subtitle=None):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.1), H)
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    txbox(slide, text, Inches(0.25), Inches(0.12), Inches(12.8), Inches(0.62),
          size=32, bold=True)
    if subtitle:
        txbox(slide, subtitle, Inches(0.25), Inches(0.75), Inches(12.8), Inches(0.38),
              size=18, italic=True, color=ACCENT)
    hline(slide, Inches(0.25), Inches(1.18), Inches(12.8), color=ACCENT)

def txbox(slide, text, x, y, w, h,
          size=22, bold=False, color=LIGHT, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = color
    return tb

def hline(slide, x, y, w, color=ACCENT, thickness=Pt(2)):
    ln = slide.shapes.add_connector(1, x, y, x + w, y)
    ln.line.color.rgb = color
    ln.line.width     = thickness

def bullet_box(slide, items, x, y, w, h, size=18, color=LIGHT, bullet="→"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        p.space_after = Pt(9)

def add_figure(slide, fname, x, y, w, h):
    path = fig(fname)
    if os.path.exists(path):
        slide.shapes.add_picture(path, x, y, w, h)
    else:
        sh = slide.shapes.add_shape(1, x, y, w, h)
        sh.fill.solid(); sh.fill.fore_color.rgb = BG2
        sh.line.color.rgb = MID
        txbox(slide, f"[missing: {fname}]",
              x + Inches(0.1), y + Inches(0.1),
              w - Inches(0.2), h - Inches(0.2), size=14, color=MID)

def caption_bar(slide, text, color=MID, size=15):
    txbox(slide, text, Inches(0.25), Inches(6.95), Inches(12.8), Inches(0.45),
          size=size, color=color, italic=True, align=PP_ALIGN.CENTER)

def code_box(slide, lines, x, y, w, h, size=16):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = BG2
    sh.line.fill.background()
    tb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1),
                                  w - Inches(0.3), h - Inches(0.2))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run(); run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = LIGHT
        run.font.name = "Courier New"
        p.space_after = Pt(4)

def finding_box(slide, tag, title, body, x, y, w=Inches(6.5), h=Inches(1.1),
                tag_color=ACC2):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = BG2
    sh.line.color.rgb = tag_color
    txbox(slide, f"{tag}  {title}", x + Inches(0.15), y + Inches(0.06),
          w - Inches(0.3), Inches(0.38), size=16, bold=True, color=tag_color)
    txbox(slide, body, x + Inches(0.15), y + Inches(0.46),
          w - Inches(0.3), Inches(0.55), size=12, color=MID)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
def slide_title(prs):
    s = blank_slide(prs); bg(s)
    bar = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.12), H)
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    txbox(s, "Certified Basin Geometry at Scale",
          Inches(0.4), Inches(0.9), Inches(12.5), Inches(1.3),
          size=44, bold=True)
    txbox(s, "How Fine-Tuning Method, Model Scale, and Training Duration\n"
             "Shape the Certified Weight-Space Ball",
          Inches(0.4), Inches(2.4), Inches(11), Inches(0.9),
          size=23, color=MID, italic=True)
    hline(s, Inches(0.4), Inches(3.5), Inches(6), color=ACCENT)

    txbox(s, "Neural Thickets  ×  Randomized Smoothing  |  Progress Update",
          Inches(0.4), Inches(3.75), Inches(9), Inches(0.55),
          size=21, color=ACC2)

    # New data summary box
    sh = s.shapes.add_shape(1, Inches(0.4), Inches(4.45), Inches(7.5), Inches(1.65))
    sh.fill.solid(); sh.fill.fore_color.rgb = BG2
    sh.line.color.rgb = GREEN
    txbox(s, "New since last update", Inches(0.55), Inches(4.52),
          Inches(7.2), Inches(0.35), size=15, bold=True, color=GREEN)
    bullet_box(s, [
        "Steps sweep (Llama-3B, lr=1e-4, 200–2000 steps) — norm grows linearly, all inside ball",
        "Trajectory run (Llama-3B, lr=2e-4, 2000 steps) — outside ball, PPL +6.4%",
        "Full FT LR sweep on Llama-3B — 3 conditions added",
        "33 total conditions, zero exceptions to the ball boundary rule",
    ], Inches(0.55), Inches(4.9), Inches(7.2), Inches(1.15), size=13, color=LIGHT, bullet="•")

    txbox(s, "Kina Kim  |  SFU  |  June 2026",
          Inches(0.4), Inches(6.55), Inches(9), Inches(0.5),
          size=17, color=MID)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Results Overview (all 4 findings)
# ══════════════════════════════════════════════════════════════════════════════
def slide_overview(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "Results Overview — 4 Findings",
              "33 conditions across 4 models, 2 LoRA ranks (r=8 and r=768), multiple LRs and step counts")

    finding_box(s, "F1", "σ½ Scaling Law",
                "σ½_pre shrinks 5× from GPT-2 (7.8e-4) → Llama-8B (1.6e-4). Better models sit in tighter basins. "
                "Consistent with Neural Thickets: denser clusters at scale → smaller certified radius.",
                Inches(0.25), Inches(1.28), tag_color=ACC2)

    finding_box(s, "F2", "Universal PPL Safety Boundary  ← MAIN RESULT",
                "norm/σ½ < 1 predicts PPL stability across 33 conditions: 4 models, 2 LoRA ranks, "
                "6 LRs, 4 step counts. Zero exceptions. The certified ball is a model-agnostic training budget.",
                Inches(0.25), Inches(2.5), tag_color=GREEN)

    finding_box(s, "F3", "Widening Depends on Scale × Rank Coverage (r/d)",
                "GPT-2 rank=8 (r/d=1%): widens 1.1–1.75×. Llama-3B rank=8 (r/d=0.26%): narrows (0.95×). "
                "Llama-3B rank=768 (r/d=25%): widens (1.18×). Widening needs sufficient r/d at a given scale.",
                Inches(0.25), Inches(3.72), tag_color=YELLOW)

    finding_box(s, "F4", "Rank Coverage Effect Reverses Across Scale",
                "GPT-2: low r/d (1%) widens, high r/d (100%) narrows (0.793×). "
                "Llama-3B: low r/d (0.26%) narrows, medium r/d (25%) widens. "
                "Direction of effect flips — scale × rank interaction, not just method.",
                Inches(0.25), Inches(4.94), tag_color=ACCENT)

    # Right column: new data
    txbox(s, "Key numbers  (new data)", Inches(7.15), Inches(1.28),
          Inches(6.0), Inches(0.4), size=17, bold=True, color=ACC2)
    code_box(s, [
        "Steps sweep (Llama-3B, lr=1e-4, rank=8):",
        "  200 steps:  norm/σ½=0.368  IN  ppl=1.000",
        "  500 steps:  norm/σ½=0.488  IN  ppl=1.000",
        " 1000 steps:  norm/σ½=0.613  IN  ppl=0.998",
        " 2000 steps:  norm/σ½=0.834  IN  ppl=0.998",
        "",
        "Trajectory (Llama-3B, lr=2e-4, 2000 steps):",
        "  norm/σ½=2.137  OUT  ppl=1.063  (+6.4%)",
        "",
        "LoRA rank=768 Llama-3B (r/d=25%, lr=1e-5):",
        "  norm/σ½=0.366  IN  ratio=1.184×  ppl=1.009",
    ], Inches(7.15), Inches(1.75), Inches(6.0), Inches(4.05), size=12)

    caption_bar(s,
        "F3/F4 reframed: 'full FT' experiments were LoRA rank=768 throughout — never true full FT.  "
        "The finding is a rank-coverage × scale interaction, not a method difference.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — F2 Main Result: Universal PPL Boundary
# ══════════════════════════════════════════════════════════════════════════════
def slide_phase_boundary(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "F2 — Universal PPL Safety Boundary  (33 conditions, zero exceptions)",
              "norm/σ½ < 1  →  PPL stable   |   norm/σ½ > 1  →  PPL degrades   |   method-agnostic")

    add_figure(s, "fig4_phase_boundary.png",
               Inches(0.25), Inches(1.28), Inches(9.0), Inches(5.3))

    txbox(s, "The rule", Inches(9.45), Inches(1.3), Inches(3.65), Inches(0.4),
          size=17, bold=True, color=GREEN)
    code_box(s, [
        "norm/σ½ < 1",
        "  → PPL ratio ≈ 1.000  ✓",
        "",
        "norm/σ½ > 1",
        "  → PPL ratio > 1.02  ✗",
    ], Inches(9.45), Inches(1.75), Inches(3.65), Inches(1.8), size=14)

    txbox(s, "Evidence", Inches(9.45), Inches(3.75), Inches(3.65), Inches(0.4),
          size=17, bold=True, color=ACC2)
    bullet_box(s, [
        "33 data points, zero exceptions",
        "● LoRA rank=8  (24 pts, 4 models × 6 LRs)",
        "▲ LoRA rank=768  (4 pts, GPT-2 r/d=100%  +  Llama-3B r/d=25%)",
        "◆ LoRA rank=8  (4 pts, steps sweep, Llama-3B)",
        "★ LoRA lr=2e-4  (1 pt, trajectory, Llama-3B)",
        "Boundary is σ½-normalized — absolute LR is irrelevant",
    ], Inches(9.45), Inches(4.2), Inches(3.65), Inches(2.5), size=12, color=MID)

    caption_bar(s,
        "New: ◆ steps sweep + ★ trajectory run added.  "
        "All 4 step counts inside ball at lr=1e-4.  "
        "Trajectory point (lr=2e-4, 2000 steps) confirmed outside ball with PPL +6.4%.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — F2 Steps Dimension (new)
# ══════════════════════════════════════════════════════════════════════════════
def slide_steps(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "F2 — Validated Across Training Duration  (new experiment)",
              "Llama-3.2-3B, SST-2, LoRA rank=8  |  lr=1e-4: 200→2000 steps  |  lr=2e-4: endpoint")

    add_figure(s, "fig10_steps_progression.png",
               Inches(0.25), Inches(1.28), Inches(9.8), Inches(5.1))

    txbox(s, "Key insight", Inches(10.25), Inches(1.3), Inches(2.85), Inches(0.4),
          size=16, bold=True, color=GREEN)
    bullet_box(s, [
        "Norm grows sub-linearly with steps (~steps^0.7 from 4 pts)",
        "At lr=1e-4: all 4 step counts inside ball, PPL stable",
        "At lr=2e-4: exits ball at step 500, PPL +6.4% at step 2000",
        "F2 holds across the time dimension",
    ], Inches(10.25), Inches(1.75), Inches(2.85), Inches(2.2), size=13, color=LIGHT)

    txbox(s, "Training budget estimate", Inches(10.25), Inches(4.15),
          Inches(2.85), Inches(0.4), size=15, bold=True, color=YELLOW)
    code_box(s, [
        "Linear model (conservative):",
        "  σ½ / norm_per_step ≈ 2400",
        "",
        "√steps model (trajectory):",
        "  (σ½ / C)²  ≈  3000",
        "",
        "True exponent uncertain",
        "  (need more step counts)",
    ], Inches(10.25), Inches(4.6), Inches(2.85), Inches(2.3), size=11)

    caption_bar(s,
        "Left: norm/σ½ vs steps — sub-linear growth (exponent ~0.7, 4 pts only).  "
        "Right: PPL stable at lr=1e-4, jumps when outside ball.  "
        "σ½_pre gives a principled budget estimate; exact formula depends on growth model.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Training Trajectory (new)
# ══════════════════════════════════════════════════════════════════════════════
def slide_trajectory(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "F2 — Basin Adapts Dynamically During Training  (new experiment)",
              "Llama-3.2-3B, LoRA rank=8, lr=2e-4  |  checkpoint every 100 steps → 20 density measurements")

    add_figure(s, "fig11_trajectory_dynamic.png",
               Inches(0.25), Inches(1.28), Inches(9.6), Inches(5.1))

    txbox(s, "What changed", Inches(10.05), Inches(1.3), Inches(3.05), Inches(0.4),
          size=15, bold=True, color=ACC2)
    bullet_box(s, [
        "Previous: single endpoint",
        "  (widening_summary only)",
        "New: 20 checkpoints,",
        "  every 100 steps",
        "Base-layer bug fixed →",
        "  all norms now valid",
    ], Inches(10.05), Inches(1.75), Inches(3.05), Inches(2.0), size=11, color=MID, bullet=" ")

    txbox(s, "Key numbers", Inches(10.05), Inches(3.85), Inches(3.05), Inches(0.4),
          size=15, bold=True, color=GREEN)
    code_box(s, [
        "Step 100:  σ½ = 1.55×",
        "  norm/σ½_ft = 0.28",
        "",
        "Step 1000: σ½ = 1.50×",
        "  norm/σ½_ft = 0.96  ✓",
        "",
        "Step 2000: σ½ = 1.40×",
        "  norm/σ½_ft = 1.21",
        "  PPL ratio = 1.063  ✓",
        "",
        "Norm ~ steps^0.51",
        "  (√steps, Adam RW)",
    ], Inches(10.05), Inches(4.3), Inches(3.05), Inches(2.55), size=11)

    caption_bar(s,
        "Left: widening ratio at each checkpoint — always > 1, jumps at step 100. "
        "Right: finetuned-ball ratio (blue) stays < 1 longer than pretrained-ball ratio (red dashed). "
        "PPL +6.3% at step 2000 even though norm/σ½_pre = 2.14.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — F1: σ½ Scaling Law
# ══════════════════════════════════════════════════════════════════════════════
def slide_scaling(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "F1 — σ½ Shrinks with Model Capability",
              "Better models sit in tighter weight-space minima  (pretrained only, no fine-tuning)")

    add_figure(s, "fig3_sigma_scaling.png",
               Inches(0.25), Inches(1.28), Inches(9.0), Inches(4.95))

    txbox(s, "Key numbers", Inches(9.5), Inches(1.3), Inches(3.6), Inches(0.4),
          size=17, bold=True, color=ACC2)
    code_box(s, [
        "GPT-2  (124M):  σ½=7.78e-4",
        "Qwen   (2.5B):  σ½=2.96e-4  (2.6×)",
        "Llama  (3.0B):  σ½=1.95e-4  (4.0×)",
        "Llama  (8.0B):  σ½=1.57e-4  (5.0×)",
    ], Inches(9.5), Inches(1.75), Inches(3.6), Inches(1.85), size=13)

    txbox(s, "Implications", Inches(9.5), Inches(3.8), Inches(3.6), Inches(0.4),
          size=17, bold=True, color=ACC2)
    bullet_box(s, [
        "σ½ tracks pretrained PPL — smaller basin = better model",
        "Consistent with Neural Thickets: denser clusters at larger scale → tighter geometry",
        "Tighter σ½ → safe LR budget shrinks at scale",
        "Thicket threshold (~1.5B) should show a kink in this curve",
    ], Inches(9.5), Inches(4.25), Inches(3.6), Inches(2.4), size=14, color=MID)

    caption_bar(s,
        "σ½_pre measured on pretrained model only.  "
        "Dashed = log-linear trend.  GPT-2 is 5× wider than Llama-3.1-8B.  "
        "Absolute LR safety threshold therefore differs by 5× across this range.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — F3/F4: Method Comparison Llama-3B
# ══════════════════════════════════════════════════════════════════════════════
def slide_method(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "F3/F4 — Widening Depends on Scale × Rank Coverage  (Llama-3B focus)",
              "Same model, same norm budget — rank coverage (r/d) determines whether basin widens or narrows")

    add_figure(s, "fig7_method_comparison.png",
               Inches(0.25), Inches(1.28), Inches(9.8), Inches(5.0))

    txbox(s, "Llama-3B results", Inches(10.2), Inches(1.3), Inches(2.9), Inches(0.4),
          size=15, bold=True, color=ACC2)
    code_box(s, [
        "LoRA rank=8:",
        "  lr=1e-4, norm/σ½=0.22",
        "  ratio=0.955×  (narrows)",
        "  PPL=1.007  ✓",
        "",
        "LoRA rank=768 (r/d=25%):",
        "  lr=1e-5, norm/σ½=0.37",
        "  ratio=1.184×  (widens!)",
        "  PPL=1.009  ✓",
        "",
        "Both inside ball →",
        "  both PPL-stable",
    ], Inches(10.2), Inches(1.75), Inches(2.9), Inches(3.2), size=11)

    txbox(s, "Interpretation", Inches(10.2), Inches(5.1), Inches(2.9), Inches(0.35),
          size=14, bold=True, color=YELLOW)
    bullet_box(s, [
        "F2 holds for both ranks",
        "rank=8 (0.26% coverage) does NOT widen at 3B — rank=768 (25%) does",
        "Higher rank LoRA accesses more directions → can reshape basin geometry",
    ], Inches(10.2), Inches(5.5), Inches(2.9), Inches(1.7), size=11, color=MID)

    caption_bar(s,
        "Left: density curves — rank=768 shifts σ½ further right despite lower norm.  "
        "Right: widening scatter — rank=768 widens (▲), rank=8 (●) narrows.  "
        "Llama-3B, SST-2, 500 steps.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — F3: Widening vs LR (4 models)
# ══════════════════════════════════════════════════════════════════════════════
def slide_widening(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "F3 — Basin Widening Depends on Model Scale",
              "GPT-2 widens inside the ball; Qwen/Llama-3B narrow; Llama-8B flat  |  LoRA rank=8")

    add_figure(s, "fig5_widening_vs_lr.png",
               Inches(0.25), Inches(1.28), Inches(9.0), Inches(5.08))

    txbox(s, "Best safe widening", Inches(9.5), Inches(1.3), Inches(3.6), Inches(0.4),
          size=16, bold=True, color=ACC2)
    code_box(s, [
        "GPT-2  (124M): 1.75×",
        "  (lr=2e-4, norm/σ½=0.21)",
        "",
        "Qwen   (2.5B): 0.82×",
        "  (lr=1e-5, norm/σ½=0.06)",
        "",
        "Llama  (3.0B): 0.99×",
        "  (lr=5e-5, norm/σ½=0.12)",
        "",
        "Llama  (8.0B): 1.03×",
        "  (lr=1e-4, norm/σ½=0.58)",
    ], Inches(9.5), Inches(1.75), Inches(3.6), Inches(3.0), size=12)

    txbox(s, "Why?", Inches(9.5), Inches(4.95), Inches(3.6), Inches(0.35),
          size=16, bold=True, color=ACC2)
    bullet_box(s, [
        "Tighter σ½ at scale → less volume to reshape inside ball",
        "LoRA rank=8 covers only ~1% of directions → basin barely moves at 3B scale",
        "Widening is a small-model phenomenon for LoRA rank=8",
    ], Inches(9.5), Inches(5.35), Inches(3.6), Inches(1.9), size=12, color=MID)

    caption_bar(s,
        "Open markers = norm/σ½ > 1 (escaped certified ball).  "
        "Dashed line = ratio=1.  lr=1e-3: catastrophic collapse (clipped for readability).")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Null Space & Subspace Analysis
# ══════════════════════════════════════════════════════════════════════════════
def slide_nullspace(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "Subspace Analysis — What Directions Does LoRA Affect?",
              "GPT-2, lr=1e-4, rank=8  |  3-way noise direction comparison")

    add_figure(s, "fig8_nullspace_comparison.png",
               Inches(0.25), Inches(1.28), Inches(9.5), Inches(5.3))

    txbox(s, "Result", Inches(9.95), Inches(1.3), Inches(3.15), Inches(0.4),
          size=17, bold=True, color=GREEN)
    code_box(s, [
        "Isotropic noise:  1.285×",
        "Subspace noise:   1.210×",
        "Null-space noise: 1.288×",
    ], Inches(9.95), Inches(1.75), Inches(3.15), Inches(1.3), size=13)

    txbox(s, "Interpretation", Inches(9.95), Inches(3.25), Inches(3.15), Inches(0.4),
          size=17, bold=True, color=ACC2)
    bullet_box(s, [
        "Null space ≈ isotropic — the other 99% of directions are barely touched",
        "Subspace slightly tighter — LoRA updates do compress the certified radius in their own directions",
        "LoRA rank=8 covers only ~1% of the full parameter space",
        "Geometric effect is concentrated in a tiny subspace",
    ], Inches(9.95), Inches(3.7), Inches(3.15), Inches(2.4), size=13, color=MID)

    caption_bar(s,
        "GPT-2 only (archive/outputs_0604).  "
        "Subspace = noise in the 8 LoRA directions.  "
        "Null space = noise in the complementary 99% of directions.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Per-Model Summary Bar
# ══════════════════════════════════════════════════════════════════════════════
def slide_summary_bar(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "Per-Model Summary at lr=1e-4  —  LoRA rank=8 vs LoRA rank=768",
              "All 4 models, SST-2, 500 steps  |  Solid=rank=8  Hatched=rank=768 (GPT-2 + Llama-3B only)")

    add_figure(s, "fig9_per_model_bar.png",
               Inches(0.25), Inches(1.28), Inches(12.83), Inches(5.22))

    caption_bar(s,
        "Left: widening ratio — GPT-2 widens, larger models flat/narrow.  "
        "Right: norm/σ½ — all below red boundary at lr=1e-4 across all 4 models.  "
        "rank=768 on GPT-2: r/d=100% (full-rank attention).  rank=768 on Llama-3B: r/d=25% (high-rank LoRA).  Qwen + Llama-8B: rank=768 not run.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE — Novelty & Triviality
# ══════════════════════════════════════════════════════════════════════════════
def slide_novelty(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "What's Actually Novel — Addressing the Triviality Concern",
              "F2 alone is vulnerable; the paper's strength is F1 + F3/F4 + the practical diagnostic")

    # ── Left: the critique + rebuttal ─────────────────────────────────────────
    txbox(s, "The critique", Inches(0.25), Inches(1.28), Inches(6.2), Inches(0.4),
          size=17, bold=True, color=ACCENT)
    bullet_box(s, [
        "σ½ is defined as the radius where performance drops",
        "F2 says: if you stay inside that radius, performance doesn't drop",
        "Isn't that just… the definition restated?",
    ], Inches(0.25), Inches(1.75), Inches(6.2), Inches(1.1), size=14, color=MID)

    txbox(s, "Why F2 is not fully trivial", Inches(0.25), Inches(2.95), Inches(6.2), Inches(0.4),
          size=17, bold=True, color=GREEN)
    bullet_box(s, [
        "σ½ is measured with isotropic Gaussian noise — random in all directions",
        "Fine-tuning is the opposite: structured, gradient-directed, low-rank",
        "No guarantee that a random-noise radius predicts robustness to a specific directional displacement",
        "Like measuring a table's wobble in all directions, then predicting whether sliding an object in one specific direction is safe",
        "That it holds across 33 conditions (4 models, 2 methods, 6 LRs, 4 step counts) is an empirical result — not a definition",
    ], Inches(0.25), Inches(3.42), Inches(6.2), Inches(2.5), size=13, color=MID)

    # vertical divider
    vbar = s.shapes.add_shape(1, Inches(6.6), Inches(1.28), Inches(0.04), Inches(5.6))
    vbar.fill.solid(); vbar.fill.fore_color.rgb = MID
    vbar.line.fill.background()

    # ── Right: what IS novel ──────────────────────────────────────────────────
    txbox(s, "What IS novel", Inches(6.8), Inches(1.28), Inches(6.3), Inches(0.4),
          size=17, bold=True, color=ACC2)

    items = [
        ("F1  σ½ Scaling Law",
         ACC2,
         "σ½ shrinks 5× from GPT-2→8B. Not definitional — measured. Better models have tighter safety budgets. Connects to Neural Thickets."),
        ("F3/F4  Rank Coverage × Scale Interaction",
         YELLOW,
         "Both conditions are LoRA — just rank=8 vs rank=768. GPT-2: low r/d widens, high r/d narrows. Llama-3B: opposite. Direction of effect reverses across scale."),
        ("Trajectory Finding",
         GREEN,
         "Basin widens within 100 steps and tracks norm throughout training. The safety envelope adapts in real time — not derivable from definitions."),
        ("σ½ as Practical Diagnostic",
         PURPLE,
         "Measurable before fine-tuning. Predicts safe LR/step combinations without running experiments. Actionable even if math is intuitive."),
    ]
    for i, (title, col, body) in enumerate(items):
        y = Inches(1.78) + Inches(1.32) * i
        sh = s.shapes.add_shape(1, Inches(6.8), y, Inches(6.3), Inches(1.2))
        sh.fill.solid(); sh.fill.fore_color.rgb = BG2
        sh.line.color.rgb = col
        txbox(s, title, Inches(6.95), y + Inches(0.06), Inches(6.0), Inches(0.38),
              size=14, bold=True, color=col)
        txbox(s, body,  Inches(6.95), y + Inches(0.46), Inches(6.0), Inches(0.65),
              size=11, color=MID)

    caption_bar(s,
        "F2 is supporting evidence (the tool works), not the main claim. "
        "The scaling law (F1) and method interaction (F3/F4) are not definitionally obvious and survive reviewer scrutiny.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Summary & Open Questions
# ══════════════════════════════════════════════════════════════════════════════
def slide_summary(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "Summary & Open Questions")

    txbox(s, "What we have", Inches(0.25), Inches(1.28),
          Inches(6.5), Inches(0.42), size=20, bold=True, color=GREEN)

    findings = [
        ("F1  σ½ Scaling Law",
         "σ½_pre shrinks 5× from GPT-2 → 8B. Better models = tighter basins. Aligns with Neural Thickets."),
        ("F2  Universal PPL Safety Boundary",
         "norm/σ½ < 1 predicts PPL stability. 33 conditions: 4 models, 2 methods, 6 LRs, 4 step counts. "
         "Zero exceptions. Temporal dimension (steps sweep) now validated."),
        ("F3  Widening Depends on Scale × Rank Coverage",
         "GPT-2 rank=8 (r/d=1%) widens 1.1–1.75×. Llama-3B rank=8 (r/d=0.26%) narrows (0.95×). "
         "Llama-3B rank=768 (r/d=25%) widens (1.18×). Sufficient r/d required at a given scale."),
        ("F4  Coverage Effect Reverses Across Scale",
         "GPT-2: low r/d (1%) widens, high r/d (100%) narrows (0.793×). "
         "Llama-3B: low r/d narrows, medium r/d widens. Scale × rank interaction — direction flips."),
    ]
    for i, (title, body) in enumerate(findings):
        y = Inches(1.75) + Inches(1.2) * i
        finding_box(s, "", title, body, Inches(0.25), y, w=Inches(6.5), tag_color=ACC2)

    vbar = s.shapes.add_shape(1, Inches(7.08), Inches(1.28), Inches(0.04), Inches(5.85))
    vbar.fill.solid(); vbar.fill.fore_color.rgb = MID
    vbar.line.fill.background()

    txbox(s, "Open Questions", Inches(7.3), Inches(1.28),
          Inches(5.8), Inches(0.42), size=20, bold=True, color=ACCENT)
    bullet_box(s, [
        "Why does GPT-2 widen but 3B LoRA narrows? Architecture, scale, pretraining diversity?",
        "Does wider σ½ after FT improve generalization, or is it a geometric artifact?",
        "Thicket threshold (~1.5B): should σ½ show a kink there? (need a 1.5B data point)",
        "Can σ½ serve as an overfitting diagnostic during training (monitor per checkpoint)?",
        "Is 'certified ball as universal LR/steps budget rule' the right paper framing?",
        "Should we pursue supervised RandOpt to compare with LoRA on a common budget?",
    ], Inches(7.3), Inches(1.78), Inches(5.8), Inches(3.8), size=14)

    hline(s, Inches(0.25), Inches(6.35), Inches(12.83), color=ACCENT, thickness=Pt(1))
    txbox(s,
          "For professor:  F3/F4 reframed — 'full FT' was LoRA rank=768 throughout, never true FT.  "
          "The finding is a rank-coverage × scale interaction. Key open question: what is the critical r/d threshold per scale?  "
          "Should we run the rank sweep (r=8→768 on Llama-3B) to find it?",
          Inches(0.25), Inches(6.48), Inches(12.83), Inches(0.85),
          size=15, bold=True, color=YELLOW)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE — Experiments Roadmap
# ══════════════════════════════════════════════════════════════════════════════
def slide_experiments(prs):
    s = blank_slide(prs); bg(s)
    title_bar(s, "Experiments Needed to Strengthen the Paper",
              "Current gaps: single task (SST-2), no true full FT, no critical rank threshold, no theoretical backing")

    # Column headers
    for x, label, col in [
        (Inches(0.25),  "Experiment",           LIGHT),
        (Inches(5.6),   "Question it answers",  ACC2),
        (Inches(9.5),   "Strengthens",          GREEN),
    ]:
        txbox(s, label, x, Inches(1.28), Inches(4.0), Inches(0.35),
              size=14, bold=True, color=col)
    hline(s, Inches(0.25), Inches(1.66), Inches(12.83), color=MID, thickness=Pt(1))

    rows = [
        # (label, priority_color, question, strengthens)
        ("① Rank sweep — Llama-3B\n   ranks 8 → 32 → 128 → 256 → 512 → 768",
         ACCENT,
         "Where is the critical coverage %\nthat flips narrowing → widening?\nIs it a fixed ratio of hidden_dim?",
         "F3/F4 become a scaling law\nnot just two data points.\nAddresses rank-triviality directly."),

        ("② Second task — AG News or MNLI\n   same models, same protocol",
         ACCENT,
         "Does F2 hold beyond binary\nsentiment? Does σ½ scaling (F1)\nreproduce on harder tasks?",
         "Reviewers will ask about\nSST-2 being too easy.\nOne task → four tasks."),

        ("③ True full fine-tuning — Llama-3B\n   all parameters, no LoRA",
         YELLOW,
         "Is F4 (method × scale) real?\nCurrent 'Full FT' is LoRA r=768\n(25% coverage), not true FT.",
         "F4 validity. Currently the\ncomparison is high-rank vs\nlow-rank LoRA, not FT vs LoRA."),

        ("④ Downstream accuracy\n   (classification F1, not just PPL)",
         YELLOW,
         "Does the boundary predict\ntask accuracy stability,\nnot just perplexity?",
         "PPL is a proxy. Accuracy is\nwhat practitioners care about.\nMakes F2 directly actionable."),

        ("⑤ 1.5B scale point\n   Llama-1B or Qwen1.5-1.5B",
         MID,
         "Is there a 'thicket threshold'\nwhere σ½ drops sharply?\nOr smooth scaling across sizes?",
         "Neural Thickets connection.\nFills gap between GPT-2 and 3B\nin the scaling law (F1)."),

        ("⑥ Quantization tolerance\n   int8/int4 post-FT",
         MID,
         "Does wider σ½ after FT predict\nbetter tolerance to quantization?\nPractical use case for σ½.",
         "Gives σ½ a concrete downstream\nuse beyond a diagnostic.\nLinks to deployment workflow."),
    ]

    row_h = Inches(0.88)
    for i, (exp, col, question, strengthens) in enumerate(rows):
        y = Inches(1.75) + row_h * i
        # Alternating row background
        if i % 2 == 0:
            bg_bar = s.shapes.add_shape(1, Inches(0.15), y - Inches(0.04),
                                        Inches(13.03), row_h - Inches(0.04))
            bg_bar.fill.solid(); bg_bar.fill.fore_color.rgb = BG2
            bg_bar.line.fill.background()

        # Priority indicator
        dot = s.shapes.add_shape(1, Inches(0.18), y + Inches(0.28), Inches(0.1), Inches(0.28))
        dot.fill.solid(); dot.fill.fore_color.rgb = col
        dot.line.fill.background()

        txbox(s, exp,        Inches(0.35), y, Inches(5.15), row_h, size=11, color=LIGHT)
        txbox(s, question,   Inches(5.6),  y, Inches(3.8),  row_h, size=11, color=MID)
        txbox(s, strengthens,Inches(9.5),  y, Inches(3.65), row_h, size=11, color=GREEN)

    hline(s, Inches(0.25), Inches(7.0), Inches(12.83), color=ACCENT, thickness=Pt(1))
    txbox(s, "● High priority  (blocks conference submission)    "
             "● Medium  (strengthens framing)    "
             "● Low  (nice to have)",
          Inches(0.25), Inches(7.08), Inches(12.83), Inches(0.35),
          size=13, color=MID, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# BUILD
# ══════════════════════════════════════════════════════════════════════════════
if __name__ == "__main__":
    prs = new_prs()
    slide_title(prs)
    slide_overview(prs)
    slide_phase_boundary(prs)
    slide_steps(prs)
    slide_trajectory(prs)
    slide_scaling(prs)
    slide_method(prs)
    slide_widening(prs)
    slide_nullspace(prs)
    slide_summary_bar(prs)
    slide_novelty(prs)
    slide_experiments(prs)
    slide_summary(prs)

    out = os.path.join(ROOT, "ntrs_presentation3.pptx")
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")
