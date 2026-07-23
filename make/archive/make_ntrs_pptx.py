"""
make_ntrs_pptx.py
=================
NTRS research PPTX: H1, H2, H3 hypotheses, experiment designs,
results + figures, verdict, and next steps.

Output: results/ntrs_research.pptx
Run:    python make_ntrs_pptx.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE = os.path.dirname(os.path.abspath(__file__))
FIGS = os.path.join(BASE, "results", "all_figures")
OUT  = os.path.join(BASE, "results", "ntrs_research.pptx")

# ── Palette ───────────────────────────────────────────────────────────────────
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1E, 0x29, 0x3B)
GRAY   = RGBColor(0x6B, 0x72, 0x80)
LGRAY  = RGBColor(0xF1, 0xF5, 0xF9)
H1C    = RGBColor(0x1D, 0x4E, 0xD8)
H2C    = RGBColor(0x05, 0x96, 0x69)
H3C    = RGBColor(0x7C, 0x3A, 0xED)
RED    = RGBColor(0xDC, 0x26, 0x26)
GREEN  = RGBColor(0x16, 0xA3, 0x4A)
AMBER  = RGBColor(0xD9, 0x77, 0x06)
LABEL  = RGBColor(0x37, 0x41, 0x51)
BLUE_L = RGBColor(0xDB, 0xEA, 0xFE)
GRN_L  = RGBColor(0xD1, 0xFA, 0xE5)
RED_L  = RGBColor(0xFE, 0xE2, 0xE2)
AMB_L  = RGBColor(0xFE, 0xF3, 0xC7)
SLATE  = RGBColor(0x94, 0xA3, 0xB8)


# ── Slide helpers ─────────────────────────────────────────────────────────────
def _prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.50)
    return prs

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _rect(slide, x, y, w, h, fill=None, line_c=None):
    shp = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line_c:
        shp.line.color.rgb = line_c
        shp.line.width = Pt(0.75)
    else:
        shp.line.width = Pt(0)
    return shp

def _tb(slide, text, x, y, w, h,
        size=13, bold=False, italic=False, color=None,
        align=PP_ALIGN.LEFT):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    if color:
        r.font.color.rgb = color
    return txb

def _bullets(slide, items, x, y, w, h, size=12, color=None):
    """items: list of str or (str, level) where level 1 = sub-bullet"""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        text, lvl = (item, 0) if isinstance(item, str) else item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = ("    ◦  " if lvl else "•  ") + text
        r.font.size = Pt(size - 1 if lvl else size)
        r.font.color.rgb = (GRAY if lvl else (color if color else LABEL))
    return txb

def _img(slide, fname, x, y, w, h=None):
    path = os.path.join(FIGS, fname) if not os.path.isabs(fname) else fname
    if not os.path.exists(path):
        print(f"  [img missing] {path}")
        return
    if h:
        slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w))

def _title(slide, text, bar_color=None, subtitle=None):
    if bar_color:
        _rect(slide, 0, 0, 13.33, 0.07, fill=bar_color)
    _tb(slide, text, 0.40, 0.12, 12.5, 0.65, size=22, bold=True, color=DARK)
    if subtitle:
        _tb(slide, subtitle, 0.40, 0.72, 12.5, 0.38, size=11, italic=True, color=GRAY)

def _insight(slide, text, bg=DARK):
    """Colored insight bar at the bottom of the slide."""
    _rect(slide, 0, 6.82, 13.33, 0.68, fill=bg)
    _tb(slide, text, 0.40, 6.86, 12.5, 0.60, size=11, color=WHITE)

def _status_badge(slide, text, x, y, color):
    _rect(slide, x, y, 2.0, 0.38, fill=color)
    _tb(slide, text, x, y, 2.0, 0.38, size=11, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER)

def _card(slide, x, y, w, h, header_color, header_text, body_items,
          body_size=11, footer_text=None, footer_color=None):
    _rect(slide, x, y, w, h, fill=LGRAY)
    _rect(slide, x, y, w, 0.50, fill=header_color)
    _tb(slide, header_text, x + 0.12, y + 0.06, w - 0.2, 0.40,
        size=12, bold=True, color=WHITE)
    _bullets(slide, body_items, x + 0.12, y + 0.58, w - 0.22,
             h - 0.80, size=body_size)
    if footer_text and footer_color:
        _rect(slide, x, y + h - 0.42, w, 0.42, fill=footer_color)
        _tb(slide, footer_text, x + 0.12, y + h - 0.40, w - 0.2, 0.38,
            size=10, bold=True, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# 1 · Title
# ══════════════════════════════════════════════════════════════════════════════
def slide_01_title(prs):
    sl = _blank(prs)
    _rect(sl, 0, 0, 13.33, 7.5, fill=DARK)
    _rect(sl, 0, 0, 0.22, 7.5, fill=H1C)
    _tb(sl, "NTRS", 0.55, 1.40, 12.0, 1.15,
        size=64, bold=True, color=WHITE)
    _tb(sl, "Neural Thickets × Randomized Smoothing",
        0.55, 2.50, 12.0, 0.70, size=22, color=SLATE)
    _tb(sl, "Research Hypotheses  ·  Experiments  ·  Findings",
        0.55, 3.18, 12.0, 0.55, size=16, color=RGBColor(0x64, 0x74, 0x8B))
    _rect(sl, 0.55, 3.88, 5.0, 0.05, fill=H1C)
    bullets = [
        ("H1   Universal Adaptation Scale",  RGBColor(0x60, 0xA5, 0xFA)),
        ("H2   Forgetting Budget",           RGBColor(0x34, 0xD3, 0x99)),
        ("H3   Model Merging",               RGBColor(0xC0, 0x84, 0xFC)),
    ]
    for i, (lbl, c) in enumerate(bullets):
        _tb(sl, lbl, 0.55, 4.10 + i * 0.58, 12.0, 0.52, size=16, color=c)
    return sl


# ══════════════════════════════════════════════════════════════════════════════
# 2 · Overview
# ══════════════════════════════════════════════════════════════════════════════
def slide_02_overview(prs):
    sl = _blank(prs)
    _title(sl, "Research Overview — Three Hypotheses Using σ½ as a Geometric Lens",
           bar_color=DARK)

    _card(sl, 0.30, 1.10, 4.15, 5.90,
          H1C, "H1  Universal Adaptation Scale",
          ["Claim: R = ‖Δθ‖/σ½_pre ≈ 1 is a universal boundary",
           "Data: 57 trajectories, 5 models, 5 tasks, 6 LR conditions each",
           "Stability boundary: varies 20× (R=0.08 to R=4.47)",
           "σ½_pre nearly constant across GPT-2 family (12× params)",
           "Llama-3B σ½ is 3–4× smaller → sharper pretrained min",
           "Redirection: σ½ = certified sharpness indicator"],
          footer_text="STATUS: REDIRECTED  (strong H1 falsified)",
          footer_color=RED)

    _card(sl, 4.68, 1.10, 4.15, 5.90,
          H2C, "H2  Forgetting Budget",
          ["Claim: ‖θ_B−θ_A‖/σ½_A > 1 predicts Task A forgetting",
           "Experiment: GPT-2 AG News→SST-2, Llama-3B SST-2→AG News",
           "Bug: --lr_a arg defined but never passed to finetune()",
           "Phase 1 accuracy = 43.9%  (expected ≥80%)",
           "σ½_A = 0.018 (30× inflated) → data unusable",
           "Fix applied: guard added, corrected commands ready"],
          footer_text="STATUS: PAUSED  (Phase 1 bug fixed, awaiting rerun)",
          footer_color=AMBER)

    _card(sl, 9.06, 1.10, 4.15, 5.90,
          H3C, "H3  Model Merging",
          ["Claim: d(θ_A, θ_B) < σ½_A + σ½_B predicts merge success",
           "Intuition: if both models within each other’s ball, the midpoint is safe",
           "Potential: geometric sufficient condition for merge success",
           "No experiments run yet",
           "Depends on H1 direction being resolved first"],
          footer_text="STATUS: DEFERRED",
          footer_color=GRAY)
    return sl


# ══════════════════════════════════════════════════════════════════════════════
# 3 · Background
# ══════════════════════════════════════════════════════════════════════════════
def slide_03_background(prs):
    sl = _blank(prs)
    _title(sl, "Foundation: Certified Density and σ½",
           bar_color=DARK,
           subtitle="The geometric quantity underlying all three hypotheses")

    _tb(sl, "Certified Density  C(θ, σ)", 0.40, 1.15, 5.8, 0.45,
        size=14, bold=True, color=DARK)
    _bullets(sl, [
        "Fraction of N=200 random Gaussian perturbations θ+ε (ε∼N(0,σ²I)) "
        "keeping NLL within eval_slack",
        "eval_slack = 1e-4 nats (WikiText-2),  1e-3 nats (classification NLL)",
        "C(θ, σ) = 1.0 → robust to all σ-scale perturbations",
        "C(θ, σ) → 0 as σ grows past the basin edge",
    ], 0.40, 1.63, 5.8, 1.80, size=12)

    _tb(sl, "σ½  (half-max sigma)", 0.40, 3.55, 5.8, 0.45,
        size=14, bold=True, color=DARK)
    _bullets(sl, [
        "σ where C(θ, σ) = 0.5 × max_density",
        "Proxy for basin width / minimum sharpness",
        "~200 forward passes per checkpoint  (cheap)",
        "Smaller σ½ = sharper minimum = better-trained model",
    ], 0.40, 4.03, 5.8, 1.80, size=12)

    _img(sl, "f2_pretrained_density.png", 6.2, 0.90, 6.85, 5.80)
    _insight(sl, "Key observation: Llama-3B σ½ = 0.00020 — 3–4× smaller than GPT-2 family (0.00060–0.00081) "
             "despite being a larger, better-trained model. σ½ tracks sharpness, not model size.")
    return sl


# ══════════════════════════════════════════════════════════════════════════════
# H1 slides
# ══════════════════════════════════════════════════════════════════════════════
def slide_04_h1_hypothesis(prs):
    sl = _blank(prs)
    _rect(sl, 0, 0, 13.33, 0.60, fill=H1C)
    _tb(sl, "H1  ·  Hypothesis", 0.40, 0.08, 6.0, 0.45,
        size=12, bold=False, color=RGBColor(0xBF, 0xDB, 0xFF))
    _tb(sl, "Universal Adaptation Scale", 0.40, 0.48, 12.5, 0.52,
        size=22, bold=True, color=WHITE)

    # Formal claim box
    _rect(sl, 0.35, 1.18, 12.6, 1.30, fill=BLUE_L)
    _tb(sl, "Claim", 0.55, 1.22, 1.8, 0.38, size=11, bold=True, color=H1C)
    _tb(sl, "R = ‖Δθ‖ / σ½_pre ≈ 1   defines a universal stability boundary",
        0.55, 1.52, 12.0, 0.52, size=17, bold=True, color=DARK)
    _tb(sl, "i.e., normalizing the update norm by σ½_pre collapses all PPL-degradation curves "
            "onto a single universal trajectory",
        0.55, 1.96, 12.0, 0.40, size=11, italic=True, color=GRAY)

    # What this means
    _tb(sl, "What this means", 0.40, 2.68, 12.5, 0.42, size=14, bold=True, color=DARK)
    _bullets(sl, [
        "R < 1 → update is inside the certified ball → pretraining quality preserved",
        "R > 1 → update exits the ball → PPL degrades",
        "R ≈ 1 holds universally across architectures (GPT-2, Llama), tasks, and LRs",
        "Testable: plot PPL ratio vs R for many (model, task, LR) conditions — should collapse",
    ], 0.40, 3.14, 12.5, 1.60, size=13)

    # Analogy box
    _rect(sl, 0.35, 4.88, 12.6, 1.65, fill=GRN_L)
    _tb(sl, "Analogy to known scaling laws", 0.55, 4.93, 5.0, 0.38,
        size=11, bold=True, color=H2C)
    _bullets(sl, [
        "Chinchilla:  N_tokens / N_params ≈ 20  collapses loss curves across model scales",
        "μP (Maximal Update Param.):  lr / width  collapses training curves across model widths",
        "H1:  ‖Δθ‖ / σ½  would collapse PPL curves across tasks, models, and LRs",
    ], 0.55, 5.32, 12.0, 1.10, size=12)
    return sl


def slide_05_h1_design(prs):
    sl = _blank(prs)
    _title(sl, "H1  ·  Experiment Design", bar_color=H1C)

    _tb(sl, "What we measure", 0.40, 1.10, 5.8, 0.42, size=13, bold=True, color=DARK)
    _bullets(sl, [
        "σ½_pre — certified density of pretrained θ  (N=200 perturbations, WikiText-2 eval)",
        "‖Δθ‖ — Frobenius norm of LoRA adapter at convergence",
        "R = ‖Δθ‖ / σ½_pre — normalized update radius",
        "WikiText-2 NLL ratio  (NLL_ft / NLL_pre) — pretraining retention",
        "Downstream task accuracy — fine-tuning success",
    ], 0.40, 1.58, 5.8, 2.20, size=12)

    _tb(sl, "What we vary", 0.40, 3.88, 5.8, 0.42, size=13, bold=True, color=DARK)
    _bullets(sl, [
        "Model:   GPT-2 (124M)  |  GPT-2-medium (354M)  |  GPT-2-large (774M)  "
        "|  GPT-2-xl (1.5B)  |  Llama-3.2-3B",
        "Task:    AG News, SST-2, MNLI, DBPedia, Yahoo Answers",
        "LR:      1e-5, 5e-5, 1e-4, 2e-4, 5e-4, 1e-3  (6 conditions per run)",
        "LoRA rank:  8 (primary)  |  rank sweep: 1–64 (GPT-2) / 8–768 (Llama-3B)",
    ], 0.40, 4.36, 5.8, 2.20, size=12)

    # Pipeline steps on right
    _tb(sl, "Pipeline (per model × task condition)", 6.55, 1.10, 6.5, 0.42,
        size=13, bold=True, color=DARK)
    steps = [
        ("1", "Load pretrained θ_pre"),
        ("2", "Measure σ½_pre  (N=200 Gaussian perturbations, WikiText-2 loss)"),
        ("3", "LoRA fine-tune  (rank=8, varying LR)"),
        ("4", "At convergence: record ‖Δθ‖,  compute R = ‖Δθ‖ / σ½_pre"),
        ("5", "Measure post-FT σ½_ft, WikiText-2 PPL ratio, task accuracy"),
        ("6", "Aggregate: plot PPL ratio vs R across all conditions"),
    ]
    y = 1.62
    for num, text in steps:
        _rect(sl, 6.55, y, 0.45, 0.44, fill=H1C)
        _tb(sl, num, 6.55, y + 0.01, 0.45, 0.40, size=14, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
        _tb(sl, text, 7.10, y, 5.8, 0.52, size=12, color=LABEL)
        y += 0.58

    _rect(sl, 6.55, 5.65, 6.5, 1.60, fill=AMB_L)
    _tb(sl, "Key constraints", 6.75, 5.70, 6.0, 0.38, size=11, bold=True, color=AMBER)
    _bullets(sl, [
        "WikiText-2 used as pretraining proxy  (not a task label)",
        "eval_slack = 1e-4 nats ensures sensitivity to small PPL changes",
        "N=200 perturbations gives ±0.04 CI at 95%  (fast, ~2 min per checkpoint)",
    ], 6.75, 6.08, 6.2, 1.10, size=11)
    return sl


def slide_06_h1_data(prs):
    sl = _blank(prs)
    _title(sl, "H1  ·  Experiments Run  +  Foundation: σ½_pre Across Models",
           bar_color=H1C)

    # Table: experiments run
    cols = ["Model", "Params", "Tasks", "LR conditions", "σ½_pre"]
    col_x = [0.32, 1.95, 3.12, 7.35, 9.35]
    col_w = [1.58, 1.12, 4.18, 1.95, 1.85]

    _rect(sl, 0.28, 1.10, 11.4, 0.46, fill=H1C)
    for h_txt, cx, cw in zip(cols, col_x, col_w):
        _tb(sl, h_txt, cx, 1.13, cw, 0.40, size=10, bold=True, color=WHITE)

    rows = [
        ("GPT-2",         "124M",  "AG News, SST-2",                     "6 each",  "0.000601"),
        ("GPT-2-medium",  "354M",  "AG News, SST-2, MNLI, DBPedia, Yahoo","6 each",  "0.000806"),
        ("GPT-2-large",   "774M",  "AG News",                            "6",       "0.000601"),
        ("GPT-2-xl",      "1.5B",  "AG News",                            "6",       "0.000665"),
        ("Llama-3.2-3B",  "3B",    "SST-2",                              "6",       "0.000200"),
    ]
    rcolors = [WHITE, LGRAY]
    for i, row in enumerate(rows):
        yr = 1.58 + i * 0.48
        _rect(sl, 0.28, yr, 11.4, 0.46, fill=rcolors[i % 2])
        for val, cx, cw in zip(row, col_x, col_w):
            _tb(sl, val, cx, yr + 0.04, cw, 0.40, size=11, color=LABEL)

    # Summary note
    _rect(sl, 0.28, 4.03, 11.4, 0.50, fill=BLUE_L)
    _tb(sl, "Total: 57 valid trajectories  (rank=8)  ·  10 (model, task) combinations  "
            "·  Note: σ½_pre nearly constant across GPT-2 family (1.3× range) despite 12× parameter scaling",
        0.42, 4.07, 11.1, 0.42, size=10, color=H1C)

    _img(sl, "f1_sigma_half_comparison.png", 0.28, 4.62, 12.6, 2.72)
    return sl


def slide_07_h1_ppl(prs):
    sl = _blank(prs)
    _title(sl, "H1  ·  Test: Does σ½ Normalize the Stability Boundary?  (WikiText-2 PPL ratio vs R)",
           bar_color=H1C,
           subtitle="If H1 holds, all curves should collapse onto one trajectory with a sharp kink at R≈1")
    _img(sl, "f3_h1_ppl_vs_R.png", 0.25, 1.20, 12.85, 5.60)
    _insight(sl,
             "Result: curves do NOT collapse. Stability boundary (where PPL first exceeds 1.01) ranges "
             "from R=0.08 (gpt2-xl/agnews) to R=4.47 (gpt2/agnews) — a 55× spread. "
             "Strong H1 is FALSIFIED. σ½ normalization is not the right normalization.")
    return sl


def slide_08_h1_acc(prs):
    sl = _blank(prs)
    _title(sl, "H1  ·  Test: Task Accuracy vs R  —  No Universal R_success",
           bar_color=H1C,
           subtitle="Different tasks require very different R values to achieve high accuracy")
    _img(sl, "f4_h1_acc_vs_R.png", 0.25, 1.20, 12.85, 5.60)
    _insight(sl,
             "gpt2-xl already knows AG News at pretraining — fine-tuning is gratuitous (high acc at tiny R). "
             "Llama-3B achieves 94.8% SST-2 inside the ball (R<1). "
             "DBPedia / Yahoo need R>1 to reach useful accuracy. "
             "No single R_success threshold predicts task success across tasks/models.")
    return sl


def slide_09_h1_boundary(prs):
    sl = _blank(prs)
    _title(sl, "H1  ·  Stability Boundary per (Model, Task)  —  20× Spread",
           bar_color=H1C,
           subtitle="R at which WikiText-2 PPL first exceeds 1.01 / 1.05 / 1.10 threshold")
    _img(sl, "f5_h1_stability_boundary.png", 0.25, 1.20, 12.85, 5.60)
    _insight(sl,
             "If H1 held, all bars would cluster near R=1. Instead: gpt2-xl/agnews goes critical at R=0.08; "
             "gpt2-medium/mnli stays stable until R=3.79. "
             "σ½ normalization partially helps (raw spread >100×, normalized spread 20×) "
             "but is far from universal collapse.")
    return sl


def slide_10_h1_verdict(prs):
    sl = _blank(prs)
    _rect(sl, 0, 0, 13.33, 0.60, fill=H1C)
    _tb(sl, "H1  ·  Verdict", 0.40, 0.08, 4.0, 0.45,
        size=12, bold=False, color=RGBColor(0xBF, 0xDB, 0xFF))
    _tb(sl, "What the Data Says", 0.40, 0.48, 12.5, 0.52,
        size=22, bold=True, color=WHITE)

    # Left: what FAILED
    _rect(sl, 0.30, 1.18, 5.85, 5.82, fill=RED_L)
    _tb(sl, "✗  What FAILED", 0.50, 1.24, 5.4, 0.45, size=14, bold=True, color=RED)
    _bullets(sl, [
        "Universal R≈1 boundary: NOT supported",
        "σ½ scales with model capacity: NOT supported",
        "Stability boundary varies 20× across tasks/models",
        "σ½_pre nearly constant across GPT-2 family (12× param scaling, <1.4× σ½ change)",
        "gpt2-xl / AG News: PPL degrades at R=0.08 (model already knew the task)",
        "gpt2-medium / MNLI: stable until R=3.79 (task is easy for pretrained model)",
    ], 0.50, 1.75, 5.5, 3.30, size=12, color=LABEL)

    # Right: what HOLDS + redirection
    _rect(sl, 6.45, 1.18, 6.58, 2.65, fill=GRN_L)
    _tb(sl, "✓  What HOLDS", 6.65, 1.24, 6.0, 0.45, size=14, bold=True, color=GREEN)
    _bullets(sl, [
        "σ½ normalization reduces 100× raw spread to 20×  (partial signal)",
        "All conditions have a PPL=1.0 safe region (some very small R)",
        "Llama-3B σ½ = 0.00020 (3–4× smaller than GPT-2) → sharper minimum",
        "Rank sweep: higher rank exits ball → forgetting accumulates",
    ], 6.65, 1.72, 6.3, 2.00, size=12, color=LABEL)

    # Redirection
    _rect(sl, 6.45, 4.00, 6.58, 3.00, fill=AMB_L)
    _tb(sl, "→  Redirection: σ½ as Certified Sharpness", 6.65, 4.06, 6.2, 0.45,
        size=13, bold=True, color=AMBER)
    _bullets(sl, [
        "σ½ ∝ 1 / sharpness of pretrained minimum",
        "Better-trained models → smaller σ½ (sharper basin) → harder to move",
        "Test: does σ½ correlate with top Hessian eigenvalue λ_max?",
        "If yes → 'certified sharpness' = concrete new contribution",
        "Practical: σ½ cheap to compute; λ_max requires power iteration",
    ], 6.65, 4.54, 6.3, 2.32, size=12, color=LABEL)
    return sl


# ══════════════════════════════════════════════════════════════════════════════
# H2 slides
# ══════════════════════════════════════════════════════════════════════════════
def slide_11_h2_hypothesis(prs):
    sl = _blank(prs)
    _rect(sl, 0, 0, 13.33, 0.60, fill=H2C)
    _tb(sl, "H2  ·  Hypothesis", 0.40, 0.08, 6.0, 0.45,
        size=12, bold=False, color=RGBColor(0xA7, 0xF3, 0xD0))
    _tb(sl, "Forgetting Budget", 0.40, 0.48, 12.5, 0.52,
        size=22, bold=True, color=WHITE)

    # Formal claim
    _rect(sl, 0.35, 1.18, 12.6, 1.25, fill=GRN_L)
    _tb(sl, "Claim", 0.55, 1.22, 1.8, 0.38, size=11, bold=True, color=H2C)
    _tb(sl, "‖θ_B − θ_A‖ / σ½_A > 1   predicts Task A forgetting after sequential A→B fine-tuning",
        0.55, 1.52, 12.0, 0.48, size=17, bold=True, color=DARK)
    _tb(sl, "i.e., whenever Phase 2 weights exit the certified ball around θ_A, Task A accuracy drops",
        0.55, 1.92, 12.0, 0.38, size=11, italic=True, color=GRAY)

    _tb(sl, "Stronger version", 0.40, 2.72, 12.5, 0.42, size=13, bold=True, color=DARK)
    _bullets(sl, [
        "Forgetting curves F(t) and A(t) align across model sizes, task pairs, and LRs",
        "The threshold is sharp at R_A = 1  (mirroring the H1 adaptation scale)",
    ], 0.40, 3.18, 12.5, 0.90, size=12)

    _tb(sl, "Why it matters", 0.40, 4.20, 12.5, 0.42, size=13, bold=True, color=DARK)
    _bullets(sl, [
        "Continual learning: how much can you train on Task B before Task A is forgotten?",
        "RLHF / SFT: σ½_A is the weight-space analog of KL penalty radius — "
        "it bounds how far RLHF policy is allowed to move from SFT policy",
        "SFT practitioners: predict forgetting BEFORE running expensive fine-tuning",
    ], 0.40, 4.66, 12.5, 1.80, size=12)
    return sl


def slide_12_h2_design(prs):
    sl = _blank(prs)
    _title(sl, "H2  ·  Experiment Design", bar_color=H2C)

    _tb(sl, "Phase 1 — Fine-tune on Task A", 0.40, 1.10, 5.8, 0.42,
        size=13, bold=True, color=DARK)
    _bullets(sl, [
        "Task A: AG News (GPT-2)  or  SST-2 (Llama-3B)",
        "LR_A = 1e-4,  steps = 2000,  rank = 8",
        "Convergence guard: raise error if Task A accuracy < 70%",
        "Save checkpoint θ_A,  measure σ½_A  (using task A classification NLL)",
    ], 0.40, 1.58, 5.8, 1.85, size=12)

    _tb(sl, "Phase 2 — Fine-tune on Task B from θ_A", 0.40, 3.55, 5.8, 0.42,
        size=13, bold=True, color=DARK)
    _bullets(sl, [
        "Task B: SST-2 (GPT-2)  or  AG News (Llama-3B)",
        "LRs = 1e-5, 1e-4, 2e-4  ×  ranks = 8, 32",
        "Every 50 steps: record R_A = ‖θ_t − θ_A‖ / σ½_A",
        "Evaluate Task A and Task B accuracy at each checkpoint",
        "Plot Task A accuracy vs R_A — drop at R_A = 1?",
    ], 0.40, 4.03, 5.8, 2.20, size=12)

    # Right: planned conditions + expected signal
    _tb(sl, "Planned task pairs", 6.55, 1.10, 6.5, 0.42, size=13, bold=True, color=DARK)
    for i, (model, pair, cmd) in enumerate([
        ("GPT-2 (124M)",    "AG News → SST-2",
         "--lr_a 1e-4  --train_steps_a 2000  --lr 1e-5 1e-4 2e-4  --ranks 8 32"),
        ("Llama-3.2-3B",    "SST-2 → AG News",
         "--lr_a 5e-5  --train_steps_a 500  --lr 1e-5 1e-4 2e-4  --ranks 8 32"),
    ]):
        y0 = 1.62 + i * 1.68
        _rect(sl, 6.55, y0, 6.5, 1.55, fill=GRN_L, line_c=H2C)
        _tb(sl, model, 6.75, y0 + 0.08, 6.0, 0.38, size=12, bold=True, color=H2C)
        _tb(sl, pair,  6.75, y0 + 0.46, 6.0, 0.38, size=13, color=DARK)
        _tb(sl, cmd,   6.75, y0 + 0.82, 6.0, 0.48, size=9,  italic=True, color=GRAY)

    _rect(sl, 6.55, 5.10, 6.5, 2.20, fill=AMB_L)
    _tb(sl, "Expected signal (if H2 holds)", 6.75, 5.16, 6.0, 0.38,
        size=11, bold=True, color=AMBER)
    _bullets(sl, [
        "Task A accuracy stays ≥80% while R_A < 1",
        "Sharp drop in Task A accuracy as R_A crosses 1",
        "If sharp: H2 confirmed  |  If gradual: H2 holds weakly",
        "If no signal: σ½_A does not capture forgetting threshold",
    ], 6.75, 5.54, 6.2, 1.60, size=11)
    return sl


def slide_13_h2_failure(prs):
    sl = _blank(prs)
    _rect(sl, 0, 0, 13.33, 0.60, fill=RED)
    _tb(sl, "H2  ·  Diagnostic", 0.40, 0.08, 6.0, 0.45,
        size=12, bold=False, color=RGBColor(0xFF, 0xCA, 0xCA))
    _tb(sl, "Phase 1 Failure — All H2 Data Is Contaminated",
        0.40, 0.48, 12.5, 0.52, size=22, bold=True, color=WHITE)
    _img(sl, "f9_seq_phase1_failure.png", 0.25, 1.18, 12.85, 5.30)
    _insight(sl,
             "Phase 1 used lr=1e-5 (default, too small). All 6 conditions: Task A accuracy = 43–46% "
             "(near-random for 4-class). σ½_A = 0.018 (30× inflated, since the checkpoint is not a "
             "converged Task A model). R_A = ‖Δθ‖ / 0.018 is meaningless. "
             "Root cause: --lr_a argument was defined in get_args() but never wired into finetune().",
             bg=RED)
    return sl


def slide_14_h2_trajectories(prs):
    sl = _blank(prs)
    _title(sl, "H2  ·  Phase 2 Forgetting Trajectories  (Despite Phase 1 Failure)",
           bar_color=H2C,
           subtitle="Task A accuracy is flat at ~44% — not forgetting, Phase 1 never worked")
    _img(sl, "f10_seq_trajectories.png", 0.25, 1.20, 12.85, 5.60)
    _insight(sl,
             "Task B (SST-2) accuracy does improve during Phase 2 (lr=1e-4, rank=32 reaches ~80%), "
             "confirming Phase 2 fine-tuning itself works. "
             "Task A (AG News) stays flat at ~44% throughout — not Phase 2 forgetting, "
             "but a Phase 1 failure artifact. R_A stays near 0 because σ½_A = 0.018 is 30× inflated.")
    return sl


def slide_15_h2_verdict(prs):
    sl = _blank(prs)
    _rect(sl, 0, 0, 13.33, 0.60, fill=H2C)
    _tb(sl, "H2  ·  Verdict", 0.40, 0.08, 4.0, 0.45,
        size=12, bold=False, color=RGBColor(0xA7, 0xF3, 0xD0))
    _tb(sl, "Current Status + Fix", 0.40, 0.48, 12.5, 0.52,
        size=22, bold=True, color=WHITE)

    _status_badge(sl, "STATUS: PAUSED", 0.40, 1.18, AMBER)
    _bullets(sl, [
        "Bug: --lr_a argument was defined in get_args() but never passed to the finetune() call",
        ("Result: Phase 1 used lr=args.lr[0] = 1e-5 (default scan start) instead of --lr_a 1e-4", 1),
        ("AG News accuracy = 43.9%  (near-random for 4-class classification)", 1),
        ("σ½_A = 0.018 (30× inflated) → R_A meaningless throughout Phase 2", 1),
    ], 0.40, 1.70, 12.5, 1.80, size=12)

    _rect(sl, 0.35, 3.60, 12.6, 1.30, fill=GRN_L)
    _tb(sl, "Fix applied (basin_widening_experiment.py)", 0.55, 3.65, 8.0, 0.38,
        size=12, bold=True, color=GREEN)
    _bullets(sl, [
        "_lr_a = args.lr_a if args.lr_a is not None else args.lr[0]   — now passed to finetune()",
        "Guard: raises RuntimeError if Phase 1 accuracy < 70%  (prevents silent contamination)",
        "Corrected command:  --lr_a 1e-4  --train_steps_a 2000",
    ], 0.55, 4.08, 12.0, 0.75, size=12)

    _rect(sl, 0.35, 5.05, 12.6, 1.00, fill=AMB_L)
    _tb(sl, "Next step: Q3 gate (cheap, ~20 min)", 0.55, 5.10, 7.0, 0.38,
        size=12, bold=True, color=AMBER)
    _tb(sl, "Run one Phase 1 check with --lr_a 1e-4 on GPT-2 / AG News. "
            "Expect accuracy ≥80%. If yes → proceed to full H2 sweep (6 LR conditions × 2 ranks). "
            "If no → increase train_steps_a or switch to Llama-3B only.",
        0.55, 5.48, 12.0, 0.52, size=11, color=LABEL)

    _insight(sl,
             "H2 is scientifically sound; the experimental infrastructure is correct. "
             "Only one argument bug prevented valid data collection. "
             "Expect 2–3 GPU-hours to produce clean H2 data.")
    return sl


# ══════════════════════════════════════════════════════════════════════════════
# 16 · H3 Hypothesis
# ══════════════════════════════════════════════════════════════════════════════
def slide_16_h3(prs):
    sl = _blank(prs)
    _rect(sl, 0, 0, 13.33, 0.60, fill=H3C)
    _tb(sl, "H3  ·  Hypothesis", 0.40, 0.08, 4.0, 0.45,
        size=12, bold=False, color=RGBColor(0xDD, 0xD6, 0xFE))
    _tb(sl, "Model Merging  —  DEFERRED", 0.40, 0.48, 12.5, 0.52,
        size=22, bold=True, color=WHITE)

    _rect(sl, 0.35, 1.18, 12.6, 1.18, fill=RGBColor(0xED, 0xE9, 0xFE))
    _tb(sl, "Claim", 0.55, 1.22, 1.8, 0.38, size=11, bold=True, color=H3C)
    _tb(sl, "d(θ_A, θ_B) < σ½_A + σ½_B   predicts that merging θ_A and θ_B preserves both tasks",
        0.55, 1.50, 12.0, 0.48, size=17, bold=True, color=DARK)
    _tb(sl, "Geometric sufficient condition: if both models are within each other’s certified ball, "
            "the midpoint lies inside both — the merged model should retain both task performances",
        0.55, 1.90, 12.0, 0.38, size=11, italic=True, color=GRAY)

    _tb(sl, "Potential experiment design", 0.40, 2.70, 12.5, 0.42, size=13, bold=True, color=DARK)
    _bullets(sl, [
        "Fine-tune GPT-2-medium separately on Task A and Task B at varying LRs",
        "Compute d(θ_A, θ_B) and compare to σ½_A + σ½_B for each (LR_A, LR_B) pair",
        "Evaluate simple-average merged model on both tasks",
        "Prediction: merge succeeds iff d < σ½_A + σ½_B",
        "Vary ranks and LRs to trace the merge boundary in (LR_A, LR_B) space",
    ], 0.40, 3.18, 12.5, 1.88, size=12)

    _rect(sl, 0.35, 5.15, 12.6, 1.00, fill=RGBColor(0xF3, 0xF4, 0xF6))
    _status_badge(sl, "STATUS: DEFERRED", 0.55, 5.22, GRAY)
    _tb(sl, "H3 requires a clean interpretation of σ½ to be meaningful. "
            "If H1 redirection confirms σ½ = certified sharpness (Q1: Hessian eigenvalue test), "
            "then H3 becomes well-motivated. Until then, the merge prediction is hard to interpret.",
        2.75, 5.22, 9.8, 0.80, size=11, color=LABEL)

    _insight(sl,
             "H3 is the most novel application of the NTRS framework — a testable geometric prediction "
             "for model merging that does NOT require running the merge. "
             "Priority: resolve H1 direction first (Q1), then H3 becomes the next major experiment.")
    return sl


# ══════════════════════════════════════════════════════════════════════════════
# 17 · Next Steps
# ══════════════════════════════════════════════════════════════════════════════
def slide_17_next_steps(prs):
    sl = _blank(prs)
    _rect(sl, 0, 0, 13.33, 7.5, fill=DARK)
    _tb(sl, "Next Steps", 0.50, 0.22, 12.3, 0.70,
        size=26, bold=True, color=WHITE)
    _rect(sl, 0.50, 0.88, 12.33, 0.04, fill=RGBColor(0x47, 0x55, 0x69))

    cards = [
        (H1C,  "Q1  [H1 core]",
         "Does σ½ correlate with top Hessian eigenvalue λ_max?",
         "Compute λ_max via power iteration on GPT-2 family.\n"
         "If σ½ ∝ 1/λ_max → 'certified sharpness' = concrete new contribution.\n"
         "Cheap: ~2 hrs per model, no fine-tuning needed.\n"
         "⭐ TOP PRIORITY",
         "YES → reframe paper around certified sharpness\n"
         "NO  → find new frame (workshop / short paper)"),

        (H2C,  "Q2  [H1 redirect]",
         "Can R_success be estimated WITHOUT fine-tuning?",
         "Approach: zero-shot gradient norm in task direction from θ_pre.\n"
         "If predictive of convergence R → full predictive theory,\n"
         "not just post-hoc description.\n"
         "Requires: task examples + one backward pass per model.",
         "YES → top-conference result (predictive, not descriptive)\n"
         "NO  → descriptive only; still useful as diagnostic"),

        (H2C,  "Q3  [H2 gate]",
         "Does Phase 1 converge cleanly with --lr_a 1e-4?",
         "Cheap: one 20-min run. GPT-2, AG News, 2000 steps.\n"
         "Guard: accuracy must reach ≥80%.\n"
         "Command:  --lr_a 1e-4  --train_steps_a 2000\n"
         "         --lr 1e-5 1e-4 2e-4  --ranks 8 32",
         "YES → full H2 sweep (~3 GPU-hrs)\n"
         "NO  → debug or switch to Llama-3B only"),

        (H3C,  "Q4  [H3 gate]",
         "Proceed to H3 Model Merging experiments?",
         "Depends on Q1: if σ½ = sharpness, H3 is well-motivated.\n"
         "Fine-tune GPT-2-medium on 2 tasks, vary LRs,\n"
         "measure d(θ_A, θ_B) vs σ½_A + σ½_B,\n"
         "evaluate simple-average merge.",
         "YES (after Q1) → novel merge prediction experiment\n"
         "NO (Q1 failed) → reframe or defer"),
    ]

    xs = [0.38, 6.85]
    for i, (color, tag, q, detail, outcome) in enumerate(cards):
        x = xs[i % 2]
        y = 1.05 + (i // 2) * 3.0
        _rect(sl, x, y, 6.15, 2.68, fill=RGBColor(0x2D, 0x3B, 0x52))
        _rect(sl, x, y, 6.15, 0.50, fill=color)
        _tb(sl, tag, x + 0.14, y + 0.06, 5.8, 0.38, size=13, bold=True, color=WHITE)
        _tb(sl, q, x + 0.14, y + 0.58, 5.8, 0.50,
            size=12, bold=True, color=RGBColor(0xE2, 0xE8, 0xF0))
        _tb(sl, detail, x + 0.14, y + 1.12, 5.8, 1.00,
            size=10, color=RGBColor(0x94, 0xA3, 0xB8))
        _tb(sl, outcome, x + 0.14, y + 2.10, 5.8, 0.52,
            size=9.5, italic=True, color=RGBColor(0x6B, 0x7A, 0x94))
    return sl


# ── Main ──────────────────────────────────────────────────────────────────────
def main():
    prs = _prs()

    slide_01_title(prs)
    slide_02_overview(prs)
    slide_03_background(prs)
    slide_04_h1_hypothesis(prs)
    slide_05_h1_design(prs)
    slide_06_h1_data(prs)
    slide_07_h1_ppl(prs)
    slide_08_h1_acc(prs)
    slide_09_h1_boundary(prs)
    slide_10_h1_verdict(prs)
    slide_11_h2_hypothesis(prs)
    slide_12_h2_design(prs)
    slide_13_h2_failure(prs)
    slide_14_h2_trajectories(prs)
    slide_15_h2_verdict(prs)
    slide_16_h3(prs)
    slide_17_next_steps(prs)

    os.makedirs(os.path.join(BASE, "results"), exist_ok=True)
    prs.save(OUT)
    print(f"Saved → {OUT}")
    print(f"  {len(prs.slides)} slides")

if __name__ == "__main__":
    main()
